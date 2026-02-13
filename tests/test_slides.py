"""Test suite for slide builders."""

from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.util import Inches

from pptx_generator.slides.registry import SLIDE_BUILDERS
from pptx_generator.config.colors import ColorPalette, PALETTES
from pptx_generator.config.typography import Typography, FontSpec, FONT_STACKS
from pptx_generator.themes.registry import get_theme


def make_prs():
    """Create a presentation with standard dimensions."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def get_test_theme_palette_typography():
    """Get standard theme, palette, and typography for testing."""
    palette = ColorPalette(**PALETTES["corporate"])
    typography = Typography(**{
        k: FontSpec(**v)
        for k, v in FONT_STACKS["professional"].items()
    })
    theme = get_theme("corporate")
    return theme, palette, typography


class TestSlideBuilderRegistry:
    """Test slide builder registry."""

    def test_slide_builders_registry_initialized(self):
        """Test that the slide builders registry is initialized."""
        assert len(SLIDE_BUILDERS) > 0

    def test_all_expected_builders_registered(self):
        """Test that all expected slide builders are registered."""
        expected_builders = [
            "title",
            "agenda",
            "section",
            "content",
            "comparison",
            "chart",
            "timeline",
            "diagram",
            "closing",
        ]

        for builder_name in expected_builders:
            assert builder_name in SLIDE_BUILDERS, f"Missing builder: {builder_name}"

    def test_all_builders_have_slide_type_property(self):
        """Test that all builders have slide_type property."""
        for name, builder in SLIDE_BUILDERS.items():
            assert hasattr(builder, "slide_type")
            assert builder.slide_type == name

    def test_all_builders_have_build_method(self):
        """Test that all builders have build method."""
        for builder in SLIDE_BUILDERS.values():
            assert hasattr(builder, "build")
            assert callable(builder.build)


class TestTitleSlideBuilder:
    """Test title slide builder."""

    def test_title_slide_basic(self):
        """Test creating a basic title slide."""
        prs = make_prs()
        theme, palette, typography = get_test_theme_palette_typography()
        builder = SLIDE_BUILDERS["title"]

        slide = builder.build(
            prs, theme, palette, typography,
            title="Test Presentation",
            subtitle="A Test Subtitle",
            author="Test Author",
            date="February 2026"
        )

        assert slide is not None
        assert len(prs.slides) == 1

    def test_title_slide_without_optional_fields(self):
        """Test title slide with minimal fields."""
        prs = make_prs()
        theme, palette, typography = get_test_theme_palette_typography()
        builder = SLIDE_BUILDERS["title"]

        slide = builder.build(
            prs, theme, palette, typography,
            title="Test Presentation"
        )

        assert slide is not None

    def test_title_slide_with_slide_number(self):
        """Test title slide with slide number."""
        prs = make_prs()
        theme, palette, typography = get_test_theme_palette_typography()
        builder = SLIDE_BUILDERS["title"]

        slide = builder.build(
            prs, theme, palette, typography,
            slide_number=1,
            title="Test Presentation",
            subtitle="Subtitle"
        )

        assert slide is not None


class TestAgendaSlideBuilder:
    """Test agenda slide builder."""

    def test_agenda_slide_basic(self):
        """Test creating a basic agenda slide."""
        prs = make_prs()
        theme, palette, typography = get_test_theme_palette_typography()
        builder = SLIDE_BUILDERS["agenda"]

        items = ["Introduction", "Key Findings", "Recommendations", "Conclusion"]
        slide = builder.build(
            prs, theme, palette, typography,
            title="Agenda",
            items=items
        )

        assert slide is not None
        assert len(prs.slides) == 1

    def test_agenda_slide_empty_items(self):
        """Test agenda slide with empty items list."""
        prs = make_prs()
        theme, palette, typography = get_test_theme_palette_typography()
        builder = SLIDE_BUILDERS["agenda"]

        slide = builder.build(
            prs, theme, palette, typography,
            title="Agenda",
            items=[]
        )

        assert slide is not None

    def test_agenda_slide_many_items(self):
        """Test agenda slide with many items."""
        prs = make_prs()
        theme, palette, typography = get_test_theme_palette_typography()
        builder = SLIDE_BUILDERS["agenda"]

        items = [f"Section {i+1}" for i in range(10)]
        slide = builder.build(
            prs, theme, palette, typography,
            title="Agenda",
            items=items
        )

        assert slide is not None


class TestSectionDividerBuilder:
    """Test section divider slide builder."""

    def test_section_slide_basic(self):
        """Test creating a basic section divider slide."""
        prs = make_prs()
        theme, palette, typography = get_test_theme_palette_typography()
        builder = SLIDE_BUILDERS["section"]

        slide = builder.build(
            prs, theme, palette, typography,
            title="Section Title",
            section_number=1
        )

        assert slide is not None

    def test_section_slide_without_number(self):
        """Test section slide without section number."""
        prs = make_prs()
        theme, palette, typography = get_test_theme_palette_typography()
        builder = SLIDE_BUILDERS["section"]

        slide = builder.build(
            prs, theme, palette, typography,
            title="Section Title"
        )

        assert slide is not None


class TestContentSlideBuilder:
    """Test content slide builder."""

    def test_content_slide_basic(self):
        """Test creating a basic content slide."""
        prs = make_prs()
        theme, palette, typography = get_test_theme_palette_typography()
        builder = SLIDE_BUILDERS["content"]

        bullets = [
            "First bullet point",
            "Second bullet point",
            "Third bullet point"
        ]

        slide = builder.build(
            prs, theme, palette, typography,
            title="Content Slide",
            bullets=bullets
        )

        assert slide is not None

    def test_content_slide_with_subtitle(self):
        """Test content slide with subtitle."""
        prs = make_prs()
        theme, palette, typography = get_test_theme_palette_typography()
        builder = SLIDE_BUILDERS["content"]

        slide = builder.build(
            prs, theme, palette, typography,
            title="Content Slide",
            subtitle="A subtitle",
            bullets=["Point 1", "Point 2"]
        )

        assert slide is not None

    def test_content_slide_no_bullets(self):
        """Test content slide without bullets."""
        prs = make_prs()
        theme, palette, typography = get_test_theme_palette_typography()
        builder = SLIDE_BUILDERS["content"]

        slide = builder.build(
            prs, theme, palette, typography,
            title="Content Slide",
            bullets=[]
        )

        assert slide is not None

    def test_content_slide_many_bullets(self):
        """Test content slide with many bullets."""
        prs = make_prs()
        theme, palette, typography = get_test_theme_palette_typography()
        builder = SLIDE_BUILDERS["content"]

        bullets = [f"Bullet point {i+1}" for i in range(8)]

        slide = builder.build(
            prs, theme, palette, typography,
            title="Content Slide",
            bullets=bullets
        )

        assert slide is not None


class TestComparisonSlideBuilder:
    """Test comparison slide builder."""

    def test_comparison_slide_basic(self):
        """Test creating a basic comparison slide."""
        prs = make_prs()
        theme, palette, typography = get_test_theme_palette_typography()
        builder = SLIDE_BUILDERS["comparison"]

        slide = builder.build(
            prs, theme, palette, typography,
            title="Comparison",
            left_title="Option A",
            left_bullets=["Pro 1", "Pro 2"],
            right_title="Option B",
            right_bullets=["Pro 1", "Pro 2"]
        )

        assert slide is not None

    def test_comparison_slide_empty_columns(self):
        """Test comparison slide with empty columns."""
        prs = make_prs()
        theme, palette, typography = get_test_theme_palette_typography()
        builder = SLIDE_BUILDERS["comparison"]

        slide = builder.build(
            prs, theme, palette, typography,
            title="Comparison",
            left_title="Left",
            left_bullets=[],
            right_title="Right",
            right_bullets=[]
        )

        assert slide is not None


class TestChartSlideBuilder:
    """Test chart slide builder."""

    def test_chart_slide_basic(self):
        """Test creating a basic chart slide."""
        prs = make_prs()
        theme, palette, typography = get_test_theme_palette_typography()
        builder = SLIDE_BUILDERS["chart"]

        chart_data = {
            "labels": ["Q1", "Q2", "Q3", "Q4"],
            "values": [100, 150, 120, 180],
            "chart_type": "bar"
        }

        slide = builder.build(
            prs, theme, palette, typography,
            title="Chart Slide",
            chart_data=chart_data
        )

        assert slide is not None

    def test_chart_slide_minimal_data(self):
        """Test chart slide with minimal data."""
        prs = make_prs()
        theme, palette, typography = get_test_theme_palette_typography()
        builder = SLIDE_BUILDERS["chart"]

        chart_data = {
            "labels": ["A"],
            "values": [10],
            "chart_type": "bar"
        }

        slide = builder.build(
            prs, theme, palette, typography,
            title="Chart Slide",
            chart_data=chart_data
        )

        assert slide is not None

    def test_chart_slide_different_chart_types(self):
        """Test chart slide with different chart types."""
        prs = make_prs()
        theme, palette, typography = get_test_theme_palette_typography()
        builder = SLIDE_BUILDERS["chart"]

        chart_types = ["bar", "line", "pie"]

        for chart_type in chart_types:
            chart_data = {
                "labels": ["A", "B", "C"],
                "values": [10, 20, 30],
                "chart_type": chart_type
            }

            slide = builder.build(
                prs, theme, palette, typography,
                title=f"{chart_type.title()} Chart",
                chart_data=chart_data
            )

            assert slide is not None


class TestTimelineSlideBuilder:
    """Test timeline slide builder."""

    def test_timeline_slide_basic(self):
        """Test creating a basic timeline slide."""
        prs = make_prs()
        theme, palette, typography = get_test_theme_palette_typography()
        builder = SLIDE_BUILDERS["timeline"]

        events = [
            {"date": "Jan 2026", "title": "Event 1"},
            {"date": "Mar 2026", "title": "Event 2"},
            {"date": "Jun 2026", "title": "Event 3"}
        ]

        slide = builder.build(
            prs, theme, palette, typography,
            title="Timeline",
            events=events
        )

        assert slide is not None

    def test_timeline_slide_empty_events(self):
        """Test timeline slide with no events."""
        prs = make_prs()
        theme, palette, typography = get_test_theme_palette_typography()
        builder = SLIDE_BUILDERS["timeline"]

        slide = builder.build(
            prs, theme, palette, typography,
            title="Timeline",
            events=[]
        )

        assert slide is not None


class TestDiagramSlideBuilder:
    """Test diagram slide builder."""

    def test_diagram_slide_basic(self):
        """Test creating a basic diagram slide."""
        prs = make_prs()
        theme, palette, typography = get_test_theme_palette_typography()
        builder = SLIDE_BUILDERS["diagram"]

        nodes = ["Step 1", "Step 2", "Step 3"]

        slide = builder.build(
            prs, theme, palette, typography,
            title="Process Diagram",
            nodes=nodes,
            diagram_type="flow"
        )

        assert slide is not None

    def test_diagram_slide_empty_nodes(self):
        """Test diagram slide with no nodes."""
        prs = make_prs()
        theme, palette, typography = get_test_theme_palette_typography()
        builder = SLIDE_BUILDERS["diagram"]

        slide = builder.build(
            prs, theme, palette, typography,
            title="Diagram",
            nodes=[],
            diagram_type="flow"
        )

        assert slide is not None

    def test_diagram_slide_different_types(self):
        """Test diagram slide with different diagram types."""
        prs = make_prs()
        theme, palette, typography = get_test_theme_palette_typography()
        builder = SLIDE_BUILDERS["diagram"]

        diagram_types = ["flow", "hierarchy", "cycle"]
        nodes = ["Node 1", "Node 2", "Node 3"]

        for diagram_type in diagram_types:
            slide = builder.build(
                prs, theme, palette, typography,
                title=f"{diagram_type.title()} Diagram",
                nodes=nodes,
                diagram_type=diagram_type
            )

            assert slide is not None


class TestClosingSlideBuilder:
    """Test closing slide builder."""

    def test_closing_slide_basic(self):
        """Test creating a basic closing slide."""
        prs = make_prs()
        theme, palette, typography = get_test_theme_palette_typography()
        builder = SLIDE_BUILDERS["closing"]

        slide = builder.build(
            prs, theme, palette, typography,
            title="Thank You",
            contact="john.doe@example.com",
            message="Questions?"
        )

        assert slide is not None

    def test_closing_slide_minimal(self):
        """Test closing slide with minimal content."""
        prs = make_prs()
        theme, palette, typography = get_test_theme_palette_typography()
        builder = SLIDE_BUILDERS["closing"]

        slide = builder.build(
            prs, theme, palette, typography,
            title="Thank You"
        )

        assert slide is not None


class TestMultipleSlides:
    """Test building multiple slides in sequence."""

    def test_build_multiple_slides(self):
        """Test building multiple slides in one presentation."""
        prs = make_prs()
        theme, palette, typography = get_test_theme_palette_typography()

        # Title slide
        SLIDE_BUILDERS["title"].build(
            prs, theme, palette, typography,
            title="Test Presentation",
            subtitle="Multi-slide Test"
        )

        # Content slide
        SLIDE_BUILDERS["content"].build(
            prs, theme, palette, typography,
            title="Content",
            bullets=["Point 1", "Point 2"]
        )

        # Closing slide
        SLIDE_BUILDERS["closing"].build(
            prs, theme, palette, typography,
            title="Thank You"
        )

        assert len(prs.slides) == 3

    def test_build_all_slide_types(self):
        """Test building one of each slide type."""
        prs = make_prs()
        theme, palette, typography = get_test_theme_palette_typography()

        # Build one of each type
        SLIDE_BUILDERS["title"].build(
            prs, theme, palette, typography,
            title="Title", subtitle="Subtitle"
        )

        SLIDE_BUILDERS["agenda"].build(
            prs, theme, palette, typography,
            title="Agenda", items=["Item 1", "Item 2"]
        )

        SLIDE_BUILDERS["section"].build(
            prs, theme, palette, typography,
            title="Section"
        )

        SLIDE_BUILDERS["content"].build(
            prs, theme, palette, typography,
            title="Content", bullets=["Point"]
        )

        SLIDE_BUILDERS["comparison"].build(
            prs, theme, palette, typography,
            title="Comparison",
            left_title="Left", left_bullets=["L1"],
            right_title="Right", right_bullets=["R1"]
        )

        SLIDE_BUILDERS["chart"].build(
            prs, theme, palette, typography,
            title="Chart",
            chart_data={"labels": ["A"], "values": [1], "chart_type": "bar"}
        )

        SLIDE_BUILDERS["timeline"].build(
            prs, theme, palette, typography,
            title="Timeline",
            events=[{"date": "2026", "title": "Event"}]
        )

        SLIDE_BUILDERS["diagram"].build(
            prs, theme, palette, typography,
            title="Diagram",
            nodes=["Node 1"], diagram_type="flow"
        )

        SLIDE_BUILDERS["closing"].build(
            prs, theme, palette, typography,
            title="Thank You"
        )

        # Should have 9 slides total
        assert len(prs.slides) == 9


class TestSlideNumbering:
    """Test slide numbering functionality."""

    def test_slide_with_number(self):
        """Test adding slide numbers."""
        prs = make_prs()
        theme, palette, typography = get_test_theme_palette_typography()

        slide = SLIDE_BUILDERS["content"].build(
            prs, theme, palette, typography,
            slide_number=1,
            title="Slide 1",
            bullets=["Content"]
        )

        assert slide is not None

    def test_slide_without_number(self):
        """Test omitting slide numbers."""
        prs = make_prs()
        theme, palette, typography = get_test_theme_palette_typography()

        slide = SLIDE_BUILDERS["content"].build(
            prs, theme, palette, typography,
            slide_number=None,
            title="Slide",
            bullets=["Content"]
        )

        assert slide is not None
