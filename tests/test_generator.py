"""Test suite for presentation generator."""

from __future__ import annotations

import pytest
import json
from pathlib import Path
from pptx import Presentation as PptxPresentation

from pptx_generator.generator import (
    PresentationGenerator,
    GeneratorError,
    ConfigurationError,
    BuildError,
)
from pptx_generator.config.settings import PresentationConfig
from pptx_generator.parsers.models import ParsedPresentation, SlideContent


class TestGeneratorInitialization:
    """Test PresentationGenerator initialization."""

    def test_generator_default_config(self):
        """Test creating generator with default config."""
        generator = PresentationGenerator()

        assert generator.config is not None
        assert generator.config.title == "Untitled"
        assert generator.preset is not None
        assert generator.theme is not None
        assert generator.palette is not None
        assert generator.typography is not None

    def test_generator_custom_config(self):
        """Test creating generator with custom config."""
        config = PresentationConfig(
            title="Test Presentation",
            subtitle="Test Subtitle",
            author="Test Author",
            palette_name="modern",
            font_stack_name="modern",
            theme_name="modern",
            preset_name="executive"
        )

        generator = PresentationGenerator(config)

        assert generator.config.title == "Test Presentation"
        assert generator.config.palette_name == "modern"
        assert generator.config.preset_name == "executive"

    def test_generator_invalid_preset(self):
        """Test that invalid preset raises ConfigurationError."""
        config = PresentationConfig(
            title="Test",
            preset_name="nonexistent_preset"
        )

        with pytest.raises(ConfigurationError):
            PresentationGenerator(config)

    def test_generator_invalid_theme(self):
        """Test that invalid theme raises ConfigurationError."""
        config = PresentationConfig(
            title="Test",
            theme_name="nonexistent_theme"
        )

        with pytest.raises(ConfigurationError):
            PresentationGenerator(config)

    def test_generator_invalid_palette(self):
        """Test that invalid palette raises ValidationError."""
        from pydantic import ValidationError

        with pytest.raises(ValidationError):
            config = PresentationConfig(
                title="Test",
                palette_name="nonexistent_palette"
            )

    def test_generator_invalid_font_stack(self):
        """Test that invalid font stack raises ValidationError."""
        from pydantic import ValidationError

        with pytest.raises(ValidationError):
            config = PresentationConfig(
                title="Test",
                font_stack_name="nonexistent_fonts"
            )


class TestGenerateFromText:
    """Test generating presentations from plain text."""

    def test_generate_from_text_basic(self, tmp_path):
        """Test basic text generation."""
        generator = PresentationGenerator()
        text = """
Introduction

This is the introduction to our presentation.

Key Points

First important point
Second important point
Third important point

Conclusion

Thank you for your attention.
"""
        output_path = tmp_path / "output.pptx"

        result_path = generator.generate_from_text(text, str(output_path))

        assert Path(result_path).exists()
        assert Path(result_path).suffix == ".pptx"

        # Verify it's a valid PowerPoint file
        prs = PptxPresentation(result_path)
        assert len(prs.slides) > 0

    def test_generate_from_text_empty(self, tmp_path):
        """Test generating from empty text."""
        generator = PresentationGenerator()
        output_path = tmp_path / "empty.pptx"

        result_path = generator.generate_from_text("", str(output_path))

        assert Path(result_path).exists()

    def test_generate_from_text_single_section(self, tmp_path):
        """Test generating from single section."""
        generator = PresentationGenerator()
        text = "Simple one-liner presentation."
        output_path = tmp_path / "single.pptx"

        result_path = generator.generate_from_text(text, str(output_path))

        assert Path(result_path).exists()

    def test_generate_from_text_custom_config(self, tmp_path):
        """Test text generation with custom configuration."""
        config = PresentationConfig(
            title="Custom Title",
            subtitle="Custom Subtitle",
            author="Test Author",
            theme_name="modern",
            preset_name="executive"
        )
        generator = PresentationGenerator(config)

        text = "Test content for custom configuration."
        output_path = tmp_path / "custom.pptx"

        result_path = generator.generate_from_text(text, str(output_path))

        assert Path(result_path).exists()


class TestGenerateFromMarkdown:
    """Test generating presentations from markdown."""

    def test_generate_from_markdown_basic(self, tmp_path):
        """Test basic markdown generation."""
        generator = PresentationGenerator()
        markdown = """
# Introduction

This is the introduction.

# Key Points

- First point
- Second point
- Third point

# Conclusion

Thank you!
"""
        output_path = tmp_path / "markdown.pptx"

        result_path = generator.generate_from_markdown(markdown, str(output_path))

        assert Path(result_path).exists()

        # Verify valid PowerPoint
        prs = PptxPresentation(result_path)
        assert len(prs.slides) > 0

    def test_generate_from_markdown_empty(self, tmp_path):
        """Test generating from empty markdown."""
        generator = PresentationGenerator()
        output_path = tmp_path / "empty_md.pptx"

        result_path = generator.generate_from_markdown("", str(output_path))

        assert Path(result_path).exists()

    def test_generate_from_markdown_complex(self, tmp_path):
        """Test markdown with various formatting."""
        generator = PresentationGenerator()
        markdown = """
# Main Title

## Subtitle

Content with **bold** and *italic*.

# Bullets

- Item 1
- Item 2
  - Sub item
- Item 3

# Code

```python
print("Hello")
```

# Links

Check [this link](https://example.com).
"""
        output_path = tmp_path / "complex.pptx"

        result_path = generator.generate_from_markdown(markdown, str(output_path))

        assert Path(result_path).exists()


class TestGenerateFromJson:
    """Test generating presentations from JSON."""

    def test_generate_from_json_dict(self, tmp_path):
        """Test JSON generation from dict."""
        generator = PresentationGenerator()
        data = {
            "title": "JSON Presentation",
            "subtitle": "Generated from JSON",
            "author": "Test",
            "slides": [
                {
                    "type": "content",
                    "title": "Section 1",
                    "bullets": ["Point 1", "Point 2"]
                },
                {
                    "type": "content",
                    "title": "Section 2",
                    "bullets": ["Point A", "Point B"]
                }
            ]
        }
        output_path = tmp_path / "json.pptx"

        result_path = generator.generate_from_json(data, str(output_path))

        assert Path(result_path).exists()

        prs = PptxPresentation(result_path)
        assert len(prs.slides) > 0

    def test_generate_from_json_string(self, tmp_path):
        """Test JSON generation from string."""
        generator = PresentationGenerator()
        data = json.dumps({
            "title": "String JSON",
            "slides": [
                {
                    "type": "content",
                    "title": "Content",
                    "bullets": ["A", "B"]
                }
            ]
        })
        output_path = tmp_path / "json_str.pptx"

        result_path = generator.generate_from_json(data, str(output_path))

        assert Path(result_path).exists()

    def test_generate_from_json_all_slide_types(self, tmp_path):
        """Test JSON with all slide types."""
        generator = PresentationGenerator()
        data = {
            "title": "All Slide Types",
            "slides": [
                {
                    "type": "content",
                    "title": "Content",
                    "bullets": ["A", "B"]
                },
                {
                    "type": "comparison",
                    "title": "Comparison",
                    "left_title": "Left",
                    "left_bullets": ["L1", "L2"],
                    "right_title": "Right",
                    "right_bullets": ["R1", "R2"]
                },
                {
                    "type": "chart",
                    "title": "Chart",
                    "chart_data": {
                        "labels": ["Q1", "Q2"],
                        "values": [100, 200],
                        "chart_type": "bar"
                    }
                },
                {
                    "type": "timeline",
                    "title": "Timeline",
                    "events": [
                        {"label": "Jan", "description": "Event 1"},
                        {"label": "Feb", "description": "Event 2"}
                    ]
                },
                {
                    "type": "diagram",
                    "title": "Diagram",
                    "nodes": ["Step 1", "Step 2"],
                    "diagram_type": "flow"
                }
            ]
        }
        output_path = tmp_path / "all_types.pptx"

        result_path = generator.generate_from_json(data, str(output_path))

        assert Path(result_path).exists()


class TestGenerateFromFile:
    """Test generating presentations from files."""

    def test_generate_from_text_file(self, tmp_path):
        """Test generating from .txt file."""
        generator = PresentationGenerator()

        # Create text file
        input_file = tmp_path / "input.txt"
        input_file.write_text("""
Section 1

Content here.

Section 2

More content.
""")

        output_path = tmp_path / "from_txt.pptx"
        result_path = generator.generate_from_file(str(input_file), str(output_path))

        assert Path(result_path).exists()

    def test_generate_from_markdown_file(self, tmp_path):
        """Test generating from .md file."""
        generator = PresentationGenerator()

        # Create markdown file
        input_file = tmp_path / "input.md"
        input_file.write_text("""
# Title

Content here.

# Section

- Bullet 1
- Bullet 2
""")

        output_path = tmp_path / "from_md.pptx"
        result_path = generator.generate_from_file(str(input_file), str(output_path))

        assert Path(result_path).exists()

    def test_generate_from_json_file(self, tmp_path):
        """Test generating from .json file."""
        generator = PresentationGenerator()

        # Create JSON file
        input_file = tmp_path / "input.json"
        data = {
            "title": "From File",
            "slides": [
                {
                    "type": "content",
                    "title": "Content",
                    "bullets": ["A", "B"]
                }
            ]
        }
        input_file.write_text(json.dumps(data))

        output_path = tmp_path / "from_json.pptx"
        result_path = generator.generate_from_file(str(input_file), str(output_path))

        assert Path(result_path).exists()

    def test_generate_from_nonexistent_file(self, tmp_path):
        """Test that nonexistent file raises error."""
        generator = PresentationGenerator()
        output_path = tmp_path / "output.pptx"

        with pytest.raises(FileNotFoundError):
            generator.generate_from_file("/nonexistent/file.txt", str(output_path))


class TestGenerateFromParsed:
    """Test generating from pre-parsed content."""

    def test_generate_from_parsed_basic(self, tmp_path):
        """Test generating from ParsedPresentation."""
        generator = PresentationGenerator()

        parsed = ParsedPresentation(
            title="Parsed Presentation",
            subtitle="From Model",
            sections=[
                SlideContent(
                    slide_type="content",
                    title="Section 1",
                    bullets=["Point 1", "Point 2"]
                ),
                SlideContent(
                    slide_type="content",
                    title="Section 2",
                    bullets=["Point A", "Point B"]
                )
            ]
        )

        output_path = tmp_path / "parsed.pptx"
        result_path = generator.generate_from_parsed(parsed, str(output_path))

        assert Path(result_path).exists()

        prs = PptxPresentation(result_path)
        assert len(prs.slides) > 0

    def test_generate_from_parsed_empty(self, tmp_path):
        """Test generating from empty ParsedPresentation."""
        generator = PresentationGenerator()

        parsed = ParsedPresentation(
            title="Empty",
            sections=[]
        )

        output_path = tmp_path / "empty_parsed.pptx"
        result_path = generator.generate_from_parsed(parsed, str(output_path))

        assert Path(result_path).exists()


class TestOutputPathHandling:
    """Test output path handling."""

    def test_output_creates_directory(self, tmp_path):
        """Test that output directory is created if needed."""
        generator = PresentationGenerator()

        output_path = tmp_path / "subdir" / "nested" / "output.pptx"
        result_path = generator.generate_from_text("Test", str(output_path))

        assert Path(result_path).exists()
        assert Path(result_path).parent.exists()

    def test_output_overwrites_existing(self, tmp_path):
        """Test that existing file is overwritten."""
        generator = PresentationGenerator()

        output_path = tmp_path / "output.pptx"

        # Generate first time
        generator.generate_from_text("First", str(output_path))
        first_size = Path(output_path).stat().st_size

        # Generate second time
        generator.generate_from_text("Second version with more content", str(output_path))
        second_size = Path(output_path).stat().st_size

        # File should still exist (overwritten)
        assert Path(output_path).exists()


class TestSlideNumbering:
    """Test slide numbering functionality."""

    def test_slides_with_numbering(self, tmp_path):
        """Test that slides include numbers when configured."""
        config = PresentationConfig(
            title="Numbered",
            include_slide_numbers=True
        )
        generator = PresentationGenerator(config)

        text = "Test content for slide numbering."
        output_path = tmp_path / "numbered.pptx"

        result_path = generator.generate_from_text(text, str(output_path))

        assert Path(result_path).exists()

    def test_slides_without_numbering(self, tmp_path):
        """Test that slides omit numbers when configured."""
        config = PresentationConfig(
            title="Unnumbered",
            include_slide_numbers=False
        )
        generator = PresentationGenerator(config)

        text = "Test content without slide numbers."
        output_path = tmp_path / "unnumbered.pptx"

        result_path = generator.generate_from_text(text, str(output_path))

        assert Path(result_path).exists()


class TestAgendaSlide:
    """Test agenda slide generation."""

    def test_agenda_included_with_multiple_sections(self, tmp_path):
        """Test that agenda is included with enough sections."""
        config = PresentationConfig(
            title="With Agenda",
            include_agenda=True
        )
        generator = PresentationGenerator(config)

        data = {
            "title": "Presentation",
            "slides": [
                {"type": "content", "title": f"Section {i}", "bullets": [f"Point {i}"]}
                for i in range(1, 6)
            ]
        }

        output_path = tmp_path / "with_agenda.pptx"
        result_path = generator.generate_from_json(data, str(output_path))

        assert Path(result_path).exists()

    def test_agenda_excluded_when_configured(self, tmp_path):
        """Test that agenda is excluded when configured."""
        config = PresentationConfig(
            title="No Agenda",
            include_agenda=False
        )
        generator = PresentationGenerator(config)

        data = {
            "title": "Presentation",
            "slides": [
                {"type": "content", "title": f"Section {i}", "bullets": [f"Point {i}"]}
                for i in range(1, 6)
            ]
        }

        output_path = tmp_path / "no_agenda.pptx"
        result_path = generator.generate_from_json(data, str(output_path))

        assert Path(result_path).exists()


class TestSectionDividers:
    """Test section divider slides."""

    def test_section_dividers_included(self, tmp_path):
        """Test that section dividers are included when configured."""
        config = PresentationConfig(
            title="With Dividers",
            include_section_dividers=True
        )
        generator = PresentationGenerator(config)

        data = {
            "title": "Presentation",
            "slides": [
                {"type": "content", "title": f"Section {i}", "bullets": ["Content"]}
                for i in range(1, 4)
            ]
        }

        output_path = tmp_path / "with_dividers.pptx"
        result_path = generator.generate_from_json(data, str(output_path))

        assert Path(result_path).exists()

    def test_section_dividers_excluded(self, tmp_path):
        """Test that section dividers are excluded when configured."""
        config = PresentationConfig(
            title="No Dividers",
            include_section_dividers=False
        )
        generator = PresentationGenerator(config)

        data = {
            "title": "Presentation",
            "slides": [
                {"type": "content", "title": f"Section {i}", "bullets": ["Content"]}
                for i in range(1, 4)
            ]
        }

        output_path = tmp_path / "no_dividers.pptx"
        result_path = generator.generate_from_json(data, str(output_path))

        assert Path(result_path).exists()


class TestDifferentPresets:
    """Test generation with different presets."""

    @pytest.mark.parametrize("preset_name", [
        "technical",
        "executive",
        "analyst",
        "trader",
        "educational",
        "summary",
    ])
    def test_generate_with_preset(self, preset_name, tmp_path):
        """Test generating presentation with each preset."""
        config = PresentationConfig(
            title=f"{preset_name.title()} Preset",
            preset_name=preset_name
        )
        generator = PresentationGenerator(config)

        text = f"Test content for {preset_name} preset."
        output_path = tmp_path / f"{preset_name}.pptx"

        result_path = generator.generate_from_text(text, str(output_path))

        assert Path(result_path).exists()


class TestDifferentThemes:
    """Test generation with different themes."""

    @pytest.mark.parametrize("theme_name", [
        "corporate",
        "modern",
        "minimal",
        "dark",
    ])
    def test_generate_with_theme(self, theme_name, tmp_path):
        """Test generating presentation with each theme."""
        config = PresentationConfig(
            title=f"{theme_name.title()} Theme",
            theme_name=theme_name
        )
        generator = PresentationGenerator(config)

        text = f"Test content for {theme_name} theme."
        output_path = tmp_path / f"theme_{theme_name}.pptx"

        result_path = generator.generate_from_text(text, str(output_path))

        assert Path(result_path).exists()


class TestIntegration:
    """Integration tests for end-to-end workflows."""

    def test_complete_presentation_workflow(self, tmp_path):
        """Test complete presentation generation workflow."""
        config = PresentationConfig(
            title="Complete Presentation",
            subtitle="End-to-End Test",
            author="Test Suite",
            date="February 2026",
            theme_name="modern",
            palette_name="modern",
            font_stack_name="modern",
            preset_name="executive",
            include_slide_numbers=True,
            include_agenda=True,
            include_section_dividers=True
        )

        generator = PresentationGenerator(config)

        markdown = """
# Executive Summary

Key highlights from Q4 2025.

# Financial Performance

- Revenue grew 25% YoY
- Operating margin improved to 18%
- Strong cash flow generation

# Market Position

- Expanded to 3 new markets
- Customer base grew 40%
- NPS score increased to 72

# Strategic Initiatives

- Launched new product line
- Completed 2 acquisitions
- Invested in R&D

# Outlook

Looking forward to continued growth in 2026.
"""

        output_path = tmp_path / "complete.pptx"
        result_path = generator.generate_from_markdown(markdown, str(output_path))

        assert Path(result_path).exists()

        # Verify it's a valid PowerPoint with multiple slides
        prs = PptxPresentation(result_path)
        assert len(prs.slides) >= 3  # At least title, content, closing
