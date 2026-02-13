"""Tests for presentation theme system."""

from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from pptx_generator.config.colors import ColorPalette
from pptx_generator.config.typography import FontSpec, Typography
from pptx_generator.themes import (
    CorporateTheme,
    DarkTheme,
    MinimalTheme,
    ModernTheme,
    get_theme,
    list_available_themes,
    register_theme,
)
from pptx_generator.themes.base import BaseTheme


@pytest.fixture
def presentation():
    """Create a test presentation."""
    return Presentation()


@pytest.fixture
def test_slide(presentation):
    """Create a test slide."""
    blank_layout = presentation.slide_layouts[6]
    return presentation.slides.add_slide(blank_layout)


@pytest.fixture
def color_palette():
    """Create a test color palette."""
    return ColorPalette(
        primary="#1E3A8A",
        secondary="#3B82F6",
        accent="#F59E0B",
        background="#FFFFFF",
        text_dark="#1F2937",
        text_light="#F9FAFB",
        success="#10B981",
        warning="#F59E0B",
        danger="#EF4444",
        neutral="#6B7280",
        gradient_start="#3B82F6",
        gradient_end="#8B5CF6",
    )


@pytest.fixture
def typography():
    """Create a test typography configuration."""
    return Typography(
        title=FontSpec(
            family="Arial", size_pt=44, bold=True, italic=False, color_hex="#1F2937"
        ),
        subtitle=FontSpec(
            family="Arial", size_pt=32, bold=False, italic=False, color_hex="#1F2937"
        ),
        heading=FontSpec(
            family="Arial", size_pt=28, bold=True, italic=False, color_hex="#1F2937"
        ),
        subheading=FontSpec(
            family="Arial", size_pt=24, bold=False, italic=False, color_hex="#1F2937"
        ),
        body=FontSpec(
            family="Arial", size_pt=18, bold=False, italic=False, color_hex="#1F2937"
        ),
        caption=FontSpec(
            family="Arial", size_pt=14, bold=False, italic=False, color_hex="#6B7280"
        ),
        footnote=FontSpec(
            family="Arial", size_pt=12, bold=False, italic=False, color_hex="#6B7280"
        ),
    )


class TestBaseTheme:
    """Tests for BaseTheme utility methods."""

    def test_hex_to_rgb_with_hash(self):
        """Test hex to RGB conversion with hash prefix."""
        rgb = BaseTheme._hex_to_rgb("#FF5733")
        expected = RGBColor(255, 87, 51)
        assert rgb == expected or (hasattr(rgb, '__iter__') and tuple(rgb) == (255, 87, 51))

    def test_hex_to_rgb_without_hash(self):
        """Test hex to RGB conversion without hash prefix."""
        rgb = BaseTheme._hex_to_rgb("FF5733")
        expected = RGBColor(255, 87, 51)
        assert rgb == expected or (hasattr(rgb, '__iter__') and tuple(rgb) == (255, 87, 51))

    def test_hex_to_rgb_black(self):
        """Test hex to RGB conversion for black."""
        rgb = BaseTheme._hex_to_rgb("#000000")
        expected = RGBColor(0, 0, 0)
        assert rgb == expected or (hasattr(rgb, '__iter__') and tuple(rgb) == (0, 0, 0))

    def test_hex_to_rgb_white(self):
        """Test hex to RGB conversion for white."""
        rgb = BaseTheme._hex_to_rgb("#FFFFFF")
        expected = RGBColor(255, 255, 255)
        assert rgb == expected or (hasattr(rgb, '__iter__') and tuple(rgb) == (255, 255, 255))

    def test_set_shape_fill(self, test_slide):
        """Test setting shape fill color."""
        shape = test_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1)
        )
        BaseTheme._set_shape_fill(shape, "#FF5733")

        assert shape.fill.type == 1  # MSO_FILL_TYPE.SOLID
        expected = RGBColor(255, 87, 51)
        assert shape.fill.fore_color.rgb == expected or tuple(shape.fill.fore_color.rgb) == (255, 87, 51)

    def test_set_shape_line(self, test_slide):
        """Test setting shape line/border."""
        shape = test_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1)
        )
        BaseTheme._set_shape_line(shape, "#0000FF", width_pt=2.5)

        expected = RGBColor(0, 0, 255)
        assert shape.line.color.rgb == expected or tuple(shape.line.color.rgb) == (0, 0, 255)

    def test_add_text_to_shape(self, test_slide):
        """Test adding styled text to shape."""
        shape = test_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(3), Inches(1)
        )
        BaseTheme._add_text_to_shape(
            shape,
            "Test Text",
            "Arial",
            18,
            "#1F2937",
            bold=True,
            italic=False,
        )

        text_frame = shape.text_frame
        assert text_frame.text == "Test Text"
        assert text_frame.paragraphs[0].runs[0].font.bold is True
        assert text_frame.paragraphs[0].runs[0].font.italic is False


class TestCorporateTheme:
    """Tests for CorporateTheme."""

    def test_name(self):
        """Test theme name property."""
        theme = CorporateTheme()
        assert theme.name == "corporate"

    def test_apply_slide_background(self, test_slide, color_palette):
        """Test applying slide background."""
        theme = CorporateTheme()
        theme.apply_slide_background(test_slide, color_palette)

        assert test_slide.background.fill.type == 1  # SOLID

    def test_add_accent_bar_left(self, test_slide, color_palette):
        """Test adding left accent bar."""
        theme = CorporateTheme()
        initial_shape_count = len(test_slide.shapes)
        theme.add_accent_bar(test_slide, color_palette, position="left")

        assert len(test_slide.shapes) == initial_shape_count + 1

    def test_add_accent_bar_positions(self, test_slide, color_palette):
        """Test accent bar in different positions."""
        theme = CorporateTheme()

        for position in ["left", "right", "top", "bottom"]:
            theme.add_accent_bar(test_slide, color_palette, position=position)

        # Should have 4 accent bars
        assert len(test_slide.shapes) == 4

    def test_add_footer(self, test_slide, color_palette, typography):
        """Test adding footer."""
        theme = CorporateTheme()
        theme.add_footer(test_slide, color_palette, typography, "Test Footer", 5)

        # Footer should add separator line + text box
        assert len(test_slide.shapes) >= 2

    def test_style_shape_primary(self, test_slide, color_palette):
        """Test styling shape with primary style."""
        theme = CorporateTheme()
        shape = test_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1)
        )

        theme.style_shape(shape, color_palette, style="primary")
        assert shape.fill.type == 1  # SOLID

    def test_style_shape_variants(self, test_slide, color_palette):
        """Test styling shape with different variants."""
        theme = CorporateTheme()

        for style in ["primary", "secondary", "accent", "neutral"]:
            shape = test_slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1)
            )
            theme.style_shape(shape, color_palette, style=style)
            assert shape.fill.type == 1  # SOLID

    def test_add_decorative_element(self, test_slide, color_palette):
        """Test adding decorative corner element."""
        theme = CorporateTheme()
        initial_count = len(test_slide.shapes)
        theme.add_decorative_element(test_slide, color_palette, element_type="corner")

        assert len(test_slide.shapes) == initial_count + 1


class TestModernTheme:
    """Tests for ModernTheme."""

    def test_name(self):
        """Test theme name property."""
        theme = ModernTheme()
        assert theme.name == "modern"

    def test_apply_slide_background(self, test_slide, color_palette):
        """Test applying slide background with header band."""
        theme = ModernTheme()
        initial_count = len(test_slide.shapes)
        theme.apply_slide_background(test_slide, color_palette)

        # Should add header band
        assert len(test_slide.shapes) == initial_count + 1

    def test_add_accent_bar_wide(self, test_slide, color_palette):
        """Test adding wide accent bar."""
        theme = ModernTheme()
        theme.add_accent_bar(test_slide, color_palette, position="left")

        assert len(test_slide.shapes) >= 1

    def test_add_footer(self, test_slide, color_palette, typography):
        """Test adding modern footer."""
        theme = ModernTheme()
        theme.add_footer(test_slide, color_palette, typography, "Modern Footer", 3)

        assert len(test_slide.shapes) >= 1

    def test_style_shape_gradient(self, test_slide, color_palette):
        """Test gradient styling for shapes."""
        theme = ModernTheme()
        shape = test_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1)
        )

        theme.style_shape(shape, color_palette, style="primary")
        assert shape.fill.type == 3  # MSO_FILL_TYPE.GRADIENT

    def test_add_decorative_geometric(self, test_slide, color_palette):
        """Test adding geometric decorative elements."""
        theme = ModernTheme()
        initial_count = len(test_slide.shapes)
        theme.add_decorative_element(test_slide, color_palette, element_type="geometric")

        # Should add multiple geometric shapes
        assert len(test_slide.shapes) > initial_count


class TestDarkTheme:
    """Tests for DarkTheme."""

    def test_name(self):
        """Test theme name property."""
        theme = DarkTheme()
        assert theme.name == "dark"

    def test_apply_slide_background(self, test_slide, color_palette):
        """Test applying dark background."""
        theme = DarkTheme()
        theme.apply_slide_background(test_slide, color_palette)

        assert test_slide.background.fill.type == 1  # SOLID

    def test_add_accent_bar_glowing(self, test_slide, color_palette):
        """Test adding glowing accent lines."""
        theme = DarkTheme()
        theme.add_accent_bar(test_slide, color_palette, position="left")

        assert len(test_slide.shapes) >= 1

    def test_add_footer_with_gradient(self, test_slide, color_palette, typography):
        """Test adding footer with gradient bar."""
        theme = DarkTheme()
        initial_count = len(test_slide.shapes)
        theme.add_footer(test_slide, color_palette, typography, "Dark Footer", 7)

        # Should add gradient bar + footer text
        assert len(test_slide.shapes) > initial_count

    def test_style_shape_transparent(self, test_slide, color_palette):
        """Test semi-transparent shape styling."""
        theme = DarkTheme()
        shape = test_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1)
        )

        theme.style_shape(shape, color_palette, style="accent")
        assert shape.fill.type == 1  # SOLID

    def test_add_decorative_element_minimal(self, test_slide, color_palette):
        """Test decorative element (minimal for dark theme)."""
        theme = DarkTheme()
        initial_count = len(test_slide.shapes)
        theme.add_decorative_element(test_slide, color_palette, element_type="gradient")

        # Dark theme has minimal decoration
        assert len(test_slide.shapes) == initial_count


class TestMinimalTheme:
    """Tests for MinimalTheme."""

    def test_name(self):
        """Test theme name property."""
        theme = MinimalTheme()
        assert theme.name == "minimal"

    def test_apply_slide_background(self, test_slide, color_palette):
        """Test applying pure white background."""
        theme = MinimalTheme()
        theme.apply_slide_background(test_slide, color_palette)

        assert test_slide.background.fill.type == 1  # SOLID
        expected = RGBColor(255, 255, 255)
        rgb = test_slide.background.fill.fore_color.rgb
        assert rgb == expected or tuple(rgb) == (255, 255, 255)

    def test_add_accent_bar_hairline(self, test_slide, color_palette):
        """Test adding thin hairline accent."""
        theme = MinimalTheme()
        theme.add_accent_bar(test_slide, color_palette, position="top")

        assert len(test_slide.shapes) >= 1

    def test_add_footer_minimal(self, test_slide, color_palette, typography):
        """Test adding minimal footer (page number only)."""
        theme = MinimalTheme()
        theme.add_footer(test_slide, color_palette, typography, "Footer", 2)

        assert len(test_slide.shapes) >= 1

    def test_add_footer_no_slide_number(self, test_slide, color_palette, typography):
        """Test footer with no slide number."""
        theme = MinimalTheme()
        initial_count = len(test_slide.shapes)
        theme.add_footer(test_slide, color_palette, typography, "", None)

        # Should not add footer if no content
        assert len(test_slide.shapes) == initial_count

    def test_style_shape_outline_only(self, test_slide, color_palette):
        """Test outline-only shape styling."""
        theme = MinimalTheme()
        shape = test_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1)
        )

        theme.style_shape(shape, color_palette, style="primary")
        assert shape.fill.type == 5  # MSO_FILL_TYPE.BACKGROUND

    def test_add_decorative_element_none(self, test_slide, color_palette):
        """Test no decorative elements."""
        theme = MinimalTheme()
        initial_count = len(test_slide.shapes)
        theme.add_decorative_element(test_slide, color_palette, element_type="none")

        # Should not add any shapes
        assert len(test_slide.shapes) == initial_count


class TestThemeRegistry:
    """Tests for theme registry."""

    def test_get_theme_corporate(self):
        """Test getting corporate theme."""
        theme = get_theme("corporate")
        assert isinstance(theme, CorporateTheme)
        assert theme.name == "corporate"

    def test_get_theme_modern(self):
        """Test getting modern theme."""
        theme = get_theme("modern")
        assert isinstance(theme, ModernTheme)
        assert theme.name == "modern"

    def test_get_theme_dark(self):
        """Test getting dark theme."""
        theme = get_theme("dark")
        assert isinstance(theme, DarkTheme)
        assert theme.name == "dark"

    def test_get_theme_minimal(self):
        """Test getting minimal theme."""
        theme = get_theme("minimal")
        assert isinstance(theme, MinimalTheme)
        assert theme.name == "minimal"

    def test_get_theme_case_insensitive(self):
        """Test theme lookup is case-insensitive."""
        theme1 = get_theme("CORPORATE")
        theme2 = get_theme("Corporate")
        theme3 = get_theme("corporate")

        assert isinstance(theme1, CorporateTheme)
        assert isinstance(theme2, CorporateTheme)
        assert isinstance(theme3, CorporateTheme)

    def test_get_theme_invalid(self):
        """Test getting invalid theme raises error."""
        with pytest.raises(ValueError, match="Theme 'invalid' not found"):
            get_theme("invalid")

    def test_list_available_themes(self):
        """Test listing available themes."""
        themes = list_available_themes()
        assert "corporate" in themes
        assert "modern" in themes
        assert "dark" in themes
        assert "minimal" in themes
        assert len(themes) == 4

    def test_register_custom_theme(self):
        """Test registering a custom theme."""

        class CustomTheme(BaseTheme):
            @property
            def name(self) -> str:
                return "custom"

            def apply_slide_background(self, slide, palette):
                pass

            def add_accent_bar(self, slide, palette, position="left"):
                pass

            def add_footer(self, slide, palette, typography, text, slide_number=None):
                pass

            def style_shape(self, shape, palette, style="primary"):
                pass

            def add_decorative_element(self, slide, palette, element_type="corner"):
                pass

        register_theme("custom", CustomTheme)
        theme = get_theme("custom")
        assert isinstance(theme, CustomTheme)

    def test_register_invalid_theme(self):
        """Test registering invalid theme raises error."""

        class NotATheme:
            pass

        with pytest.raises(ValueError, match="must inherit from BaseTheme"):
            register_theme("invalid", NotATheme)


class TestThemeIntegration:
    """Integration tests for complete theme application."""

    def test_full_corporate_theme_application(self, test_slide, color_palette, typography):
        """Test applying all corporate theme elements."""
        theme = CorporateTheme()

        theme.apply_slide_background(test_slide, color_palette)
        theme.add_accent_bar(test_slide, color_palette, position="left")
        theme.add_footer(test_slide, color_palette, typography, "Test Presentation", 1)
        theme.add_decorative_element(test_slide, color_palette, element_type="corner")

        # Add a styled shape
        shape = test_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(2), Inches(2), Inches(3), Inches(1.5)
        )
        theme.style_shape(shape, color_palette, style="primary")

        # Verify all elements were added
        assert len(test_slide.shapes) >= 4

    def test_full_modern_theme_application(self, test_slide, color_palette, typography):
        """Test applying all modern theme elements."""
        theme = ModernTheme()

        theme.apply_slide_background(test_slide, color_palette)
        theme.add_accent_bar(test_slide, color_palette, position="left")
        theme.add_footer(test_slide, color_palette, typography, "Modern Deck", 2)
        theme.add_decorative_element(test_slide, color_palette, element_type="geometric")

        # Add a styled shape with gradient
        shape = test_slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(2), Inches(3), Inches(1.5)
        )
        theme.style_shape(shape, color_palette, style="accent")

        assert len(test_slide.shapes) >= 5

    def test_theme_switching(self, presentation, color_palette, typography):
        """Test switching between different themes on different slides."""
        blank_layout = presentation.slide_layouts[6]

        # Corporate slide
        slide1 = presentation.slides.add_slide(blank_layout)
        corporate = CorporateTheme()
        corporate.apply_slide_background(slide1, color_palette)
        corporate.add_accent_bar(slide1, color_palette)

        # Modern slide
        slide2 = presentation.slides.add_slide(blank_layout)
        modern = ModernTheme()
        modern.apply_slide_background(slide2, color_palette)
        modern.add_accent_bar(slide2, color_palette)

        # Dark slide
        slide3 = presentation.slides.add_slide(blank_layout)
        dark = DarkTheme()
        dark.apply_slide_background(slide3, color_palette)
        dark.add_accent_bar(slide3, color_palette)

        # Minimal slide
        slide4 = presentation.slides.add_slide(blank_layout)
        minimal = MinimalTheme()
        minimal.apply_slide_background(slide4, color_palette)
        minimal.add_accent_bar(slide4, color_palette)

        assert len(presentation.slides) == 4
