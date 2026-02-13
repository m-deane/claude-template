"""Main presentation generator orchestrator.

This module provides the PresentationGenerator class, which coordinates all
components (parsers, themes, presets, slide builders) to generate PowerPoint
presentations from various input formats.
"""

from __future__ import annotations

import os
from pathlib import Path
from datetime import datetime
import logging
from typing import TYPE_CHECKING

from pptx import Presentation
from pptx.util import Inches

from pptx_generator.config.colors import ColorPalette, PALETTES
from pptx_generator.config.typography import Typography, FontSpec, FONT_STACKS
from pptx_generator.config.settings import PresentationConfig
from pptx_generator.themes.registry import get_theme
from pptx_generator.presets.registry import get_preset
from pptx_generator.slides.registry import SLIDE_BUILDERS
from pptx_generator.parsers.models import ParsedPresentation, SlideContent
from pptx_generator.parsers import TextParser, MarkdownParser, JsonParser

if TYPE_CHECKING:
    from pptx_generator.themes.base import BaseTheme
    from pptx_generator.presets.base import BasePreset

logger = logging.getLogger(__name__)


class GeneratorError(Exception):
    """Base exception for generator errors."""
    pass


class ConfigurationError(GeneratorError):
    """Raised when configuration is invalid."""
    pass


class BuildError(GeneratorError):
    """Raised when slide building fails."""
    pass


class PresentationGenerator:
    """Main orchestrator for generating PowerPoint presentations.

    This class ties together all components of the presentation generation
    pipeline: parsers, themes, presets, palettes, typography, and slide
    builders.

    Attributes:
        config: Complete presentation configuration
        preset: Active preset instance
        theme: Active theme instance
        palette: Active color palette
        typography: Active typography system

    Example:
        >>> from pptx_generator.config.settings import PresentationConfig
        >>> config = PresentationConfig(
        ...     title="Q4 Results",
        ...     subtitle="Financial Overview",
        ...     author="Jane Smith",
        ...     theme_name="corporate",
        ...     preset_name="executive"
        ... )
        >>> generator = PresentationGenerator(config)
        >>> output_path = generator.generate_from_text(
        ...     "Our revenue increased...",
        ...     output_path="output/q4_results.pptx"
        ... )
    """

    def __init__(self, config: PresentationConfig | None = None) -> None:
        """Initialize the presentation generator.

        Args:
            config: Presentation configuration. If None, uses default config
                   with title "Untitled".

        Raises:
            ConfigurationError: If configuration is invalid or references
                              non-existent themes/presets/palettes.
        """
        self.config = config or PresentationConfig(title="Untitled")
        try:
            self._setup()
        except (ValueError, KeyError) as e:
            raise ConfigurationError(f"Failed to initialize generator: {e}") from e

    def _setup(self) -> None:
        """Initialize theme, palette, typography, and preset from config.

        Raises:
            ValueError: If theme, preset, palette, or font stack names are invalid.
        """
        # Get preset - it may provide recommendations for theme/palette/fonts
        try:
            self.preset: BasePreset = get_preset(self.config.preset_name)
        except ValueError as e:
            logger.error(f"Failed to load preset '{self.config.preset_name}': {e}")
            raise

        # Use config values - they're already validated by PresentationConfig
        theme_name = self.config.theme_name
        palette_name = self.config.palette_name
        font_stack_name = self.config.font_stack_name

        # Load theme
        try:
            self.theme: BaseTheme = get_theme(theme_name)
        except ValueError as e:
            logger.error(f"Failed to load theme '{theme_name}': {e}")
            raise

        # Load palette
        if palette_name not in PALETTES:
            available = ', '.join(sorted(PALETTES.keys()))
            raise ValueError(
                f"Palette '{palette_name}' not found. Available: {available}"
            )
        self.palette = ColorPalette(**PALETTES[palette_name])

        # Load typography
        if font_stack_name not in FONT_STACKS:
            available = ', '.join(sorted(FONT_STACKS.keys()))
            raise ValueError(
                f"Font stack '{font_stack_name}' not found. Available: {available}"
            )

        # Build Typography instance from font stack dictionary
        font_stack = FONT_STACKS[font_stack_name]
        self.typography = Typography(**{
            level: FontSpec(**font_stack[level])
            for level in ["title", "subtitle", "heading", "subheading",
                         "body", "caption", "footnote"]
        })

        logger.info(
            f"Generator initialized: theme={theme_name}, "
            f"palette={palette_name}, font_stack={font_stack_name}, "
            f"preset={self.config.preset_name}"
        )

    def generate_from_text(
        self,
        text: str,
        output_path: str = "output/presentation.pptx"
    ) -> str:
        """Generate presentation from plain text content.

        Args:
            text: Plain text content with sections separated by blank lines
            output_path: Path where PowerPoint file will be saved

        Returns:
            Absolute path to the generated PowerPoint file

        Raises:
            GeneratorError: If parsing or generation fails
            IOError: If output file cannot be written

        Example:
            >>> generator = PresentationGenerator()
            >>> text = '''
            ... Introduction
            ... This is our quarterly review.
            ...
            ... Key Metrics
            ... Revenue increased by 25%
            ... User growth exceeded targets
            ... '''
            >>> path = generator.generate_from_text(text)
        """
        try:
            parser = TextParser()
            parsed = parser.parse(
                text,
                max_bullets=self.preset.max_bullets_per_slide,
                max_words_per_bullet=self.preset.max_words_per_bullet
            )

            # Override parsed title with config title if explicitly set
            if self.config.title != "Untitled":
                parsed.title = self.config.title
            if self.config.subtitle:
                parsed.subtitle = self.config.subtitle

            return self._generate(parsed, output_path)

        except Exception as e:
            logger.error(f"Failed to generate from text: {e}")
            raise GeneratorError(f"Text generation failed: {e}") from e

    def generate_from_markdown(
        self,
        markdown_text: str,
        output_path: str = "output/presentation.pptx"
    ) -> str:
        """Generate presentation from markdown content.

        Supports standard markdown syntax:
        - # Headings for slide titles
        - * or - for bullet points
        - ## Subheadings
        - Code blocks, links, emphasis

        Args:
            markdown_text: Markdown formatted content
            output_path: Path where PowerPoint file will be saved

        Returns:
            Absolute path to the generated PowerPoint file

        Raises:
            GeneratorError: If parsing or generation fails
            IOError: If output file cannot be written

        Example:
            >>> markdown = '''
            ... # Introduction
            ... Welcome to our presentation
            ...
            ... # Key Points
            ... - First point
            ... - Second point
            ... '''
            >>> path = generator.generate_from_markdown(markdown)
        """
        try:
            parser = MarkdownParser()
            parsed = parser.parse(
                markdown_text,
                max_bullets=self.preset.max_bullets_per_slide
            )

            if self.config.title != "Untitled":
                parsed.title = self.config.title
            if self.config.subtitle:
                parsed.subtitle = self.config.subtitle

            return self._generate(parsed, output_path)

        except Exception as e:
            logger.error(f"Failed to generate from markdown: {e}")
            raise GeneratorError(f"Markdown generation failed: {e}") from e

    def generate_from_json(
        self,
        json_data: dict | str,
        output_path: str = "output/presentation.pptx"
    ) -> str:
        """Generate presentation from structured JSON data.

        JSON format should include:
        - title, subtitle, author, date (metadata)
        - sections: list of slide objects with type, title, content

        Args:
            json_data: JSON string or dict containing presentation structure
            output_path: Path where PowerPoint file will be saved

        Returns:
            Absolute path to the generated PowerPoint file

        Raises:
            GeneratorError: If parsing or generation fails
            IOError: If output file cannot be written

        Example:
            >>> json_data = {
            ...     "title": "My Presentation",
            ...     "sections": [
            ...         {
            ...             "slide_type": "content",
            ...             "title": "Introduction",
            ...             "bullets": ["Point 1", "Point 2"]
            ...         }
            ...     ]
            ... }
            >>> path = generator.generate_from_json(json_data)
        """
        try:
            parser = JsonParser()
            parsed = parser.parse(json_data)
            return self._generate(parsed, output_path)

        except Exception as e:
            logger.error(f"Failed to generate from JSON: {e}")
            raise GeneratorError(f"JSON generation failed: {e}") from e

    def generate_from_file(
        self,
        file_path: str,
        output_path: str = "output/presentation.pptx"
    ) -> str:
        """Auto-detect file type and generate presentation.

        Supported file types:
        - .json: Structured JSON presentation data
        - .md, .markdown: Markdown formatted content
        - .txt or other: Plain text content

        Args:
            file_path: Path to input file
            output_path: Path where PowerPoint file will be saved

        Returns:
            Absolute path to the generated PowerPoint file

        Raises:
            FileNotFoundError: If input file doesn't exist
            GeneratorError: If parsing or generation fails
            IOError: If output file cannot be written

        Example:
            >>> path = generator.generate_from_file("input/content.md")
        """
        path = Path(file_path)

        if not path.exists():
            raise FileNotFoundError(f"Input file not found: {file_path}")

        suffix = path.suffix.lower()

        try:
            if suffix == ".json":
                parser = JsonParser()
                parsed = parser.parse_file(file_path)
            elif suffix in (".md", ".markdown"):
                parser = MarkdownParser()
                parsed = parser.parse_file(
                    file_path,
                    max_bullets=self.preset.max_bullets_per_slide
                )
            else:  # .txt or any other
                parser = TextParser()
                parsed = parser.parse_file(
                    file_path,
                    max_bullets=self.preset.max_bullets_per_slide,
                    max_words_per_bullet=self.preset.max_words_per_bullet
                )

            # Override with config metadata if explicitly set
            if self.config.title != "Untitled":
                parsed.title = self.config.title
            if self.config.subtitle:
                parsed.subtitle = self.config.subtitle

            return self._generate(parsed, output_path)

        except Exception as e:
            logger.error(f"Failed to generate from file '{file_path}': {e}")
            raise GeneratorError(f"File generation failed: {e}") from e

    def generate_from_parsed(
        self,
        parsed: ParsedPresentation,
        output_path: str = "output/presentation.pptx"
    ) -> str:
        """Generate presentation from pre-parsed content.

        Use this method when you have already parsed content into a
        ParsedPresentation object and want full control over the data.

        Args:
            parsed: Pre-parsed presentation data
            output_path: Path where PowerPoint file will be saved

        Returns:
            Absolute path to the generated PowerPoint file

        Raises:
            GeneratorError: If generation fails
            IOError: If output file cannot be written

        Example:
            >>> from pptx_generator.parsers.models import (
            ...     ParsedPresentation, SlideContent
            ... )
            >>> parsed = ParsedPresentation(
            ...     title="Custom Presentation",
            ...     sections=[
            ...         SlideContent(
            ...             slide_type="content",
            ...             title="Section 1",
            ...             bullets=["Point A", "Point B"]
            ...         )
            ...     ]
            ... )
            >>> path = generator.generate_from_parsed(parsed)
        """
        try:
            return self._generate(parsed, output_path)
        except Exception as e:
            logger.error(f"Failed to generate from parsed data: {e}")
            raise GeneratorError(f"Generation from parsed data failed: {e}") from e

    def _generate(self, parsed: ParsedPresentation, output_path: str) -> str:
        """Core generation logic.

        Args:
            parsed: Parsed presentation content
            output_path: Path where PowerPoint file will be saved

        Returns:
            Absolute path to the generated PowerPoint file

        Raises:
            BuildError: If slide building fails
            IOError: If file cannot be saved
        """
        logger.info(f"Starting presentation generation: {parsed.title}")

        # Create presentation with correct dimensions
        prs = Presentation()
        prs.slide_width = Inches(self.config.slide_config.width_inches)
        prs.slide_height = Inches(self.config.slide_config.height_inches)

        # Build the slide sequence based on preset and parsed content
        slides_to_build = self._plan_slides(parsed)

        logger.info(f"Building {len(slides_to_build)} slides")

        # Build each slide
        for i, slide_content in enumerate(slides_to_build):
            try:
                slide_num = (i + 1) if self.config.include_slide_numbers else None
                self._build_slide(prs, slide_content, slide_num)
            except Exception as e:
                logger.error(
                    f"Failed to build slide {i + 1} (type={slide_content.slide_type}): {e}"
                )
                raise BuildError(
                    f"Failed to build slide {i + 1} "
                    f"(type={slide_content.slide_type}): {e}"
                ) from e

        # Ensure output directory exists
        output_abs = os.path.abspath(output_path)
        output_dir = os.path.dirname(output_abs)

        if output_dir:
            os.makedirs(output_dir, exist_ok=True)

        # Save presentation
        try:
            prs.save(output_abs)
            logger.info(f"Presentation saved successfully: {output_abs}")
        except Exception as e:
            logger.error(f"Failed to save presentation: {e}")
            raise IOError(f"Failed to save presentation to '{output_abs}': {e}") from e

        return output_abs

    def _plan_slides(self, parsed: ParsedPresentation) -> list[SlideContent]:
        """Plan the slide sequence using preset rules and parsed content.

        This method determines which slides to include based on:
        - Preset flags (use_title_slide, use_agenda_slide, etc.)
        - Config options (include_agenda, include_section_dividers)
        - Parsed content structure

        Args:
            parsed: Parsed presentation content

        Returns:
            List of SlideContent objects in presentation order
        """
        slides: list[SlideContent] = []

        # Title slide
        if self.preset.use_title_slide:
            slides.append(SlideContent(
                slide_type="title",
                title=parsed.title or self.config.title,
                subtitle=parsed.subtitle or self.config.subtitle,
            ))

        # Collect section titles for agenda
        section_titles: list[str] = []
        for section in parsed.sections:
            if section.title and section.title not in section_titles:
                # Only add unique, non-empty titles
                section_titles.append(section.title)

        # Agenda slide - only if we have multiple sections
        if (self.preset.use_agenda_slide and
            self.config.include_agenda and
            len(section_titles) > 2):
            slides.append(SlideContent(
                slide_type="agenda",
                title="Agenda",
                bullets=section_titles[:10],  # Cap at 10 items for readability
            ))

        # Content slides from parsed sections
        section_num = 0
        for i, section in enumerate(parsed.sections):
            # Track section numbering for content slides
            if section.slide_type == "content":
                section_num += 1

                # Add section divider before content (not for first section)
                if (self.preset.use_section_dividers and
                    self.config.include_section_dividers and
                    section_num > 1 and
                    section.title):
                    slides.append(SlideContent(
                        slide_type="section",
                        title=section.title,
                        section_number=section_num,
                    ))

            # Add the actual content slide
            slides.append(section)

        # Closing slide
        if self.preset.use_closing_slide:
            slides.append(SlideContent(
                slide_type="closing",
                title="Thank You",
                contact=self.config.author or "",
                message=parsed.subtitle or "",
            ))

        logger.debug(f"Planned {len(slides)} slides")
        return slides

    def _build_slide(
        self,
        prs: Presentation,
        slide_content: SlideContent,
        slide_num: int | None
    ) -> None:
        """Build a single slide using the appropriate builder.

        Args:
            prs: Python-pptx Presentation object
            slide_content: Content for this slide
            slide_num: Slide number (1-indexed) or None to omit numbering

        Raises:
            BuildError: If slide builder fails
        """
        builder = SLIDE_BUILDERS.get(slide_content.slide_type)

        if builder is None:
            # Fall back to content builder for unknown types
            logger.warning(
                f"No builder found for slide type '{slide_content.slide_type}', "
                f"falling back to 'content' builder"
            )
            builder = SLIDE_BUILDERS.get("content")

            if builder is None:
                raise BuildError(
                    f"No builder available for slide type '{slide_content.slide_type}' "
                    f"and no fallback 'content' builder found"
                )

        # Build kwargs from SlideContent
        kwargs = self._build_kwargs(slide_content)

        # Build the slide
        logger.debug(
            f"Building slide: type={slide_content.slide_type}, "
            f"title={slide_content.title}, num={slide_num}"
        )

        builder.build(
            prs,
            self.theme,
            self.palette,
            self.typography,
            slide_number=slide_num,
            **kwargs
        )

    def _build_kwargs(self, sc: SlideContent) -> dict:
        """Convert SlideContent to kwargs dict for slide builders.

        Args:
            sc: SlideContent object

        Returns:
            Dictionary of kwargs appropriate for the slide type's builder
        """
        kwargs: dict = {}

        if sc.slide_type == "title":
            kwargs = {
                "title": sc.title,
                "subtitle": sc.subtitle,
                "author": self.config.author,
                "date": self.config.date or datetime.now().strftime("%B %Y")
            }
        elif sc.slide_type == "agenda":
            kwargs = {
                "title": sc.title or "Agenda",
                "items": sc.bullets
            }
        elif sc.slide_type == "section":
            kwargs = {
                "title": sc.title,
                "section_number": sc.section_number
            }
        elif sc.slide_type == "content":
            kwargs = {
                "title": sc.title,
                "bullets": sc.bullets,
                "subtitle": sc.subtitle
            }
        elif sc.slide_type == "comparison":
            kwargs = {
                "title": sc.title,
                "left_title": sc.left_title,
                "left_bullets": sc.left_bullets,
                "right_title": sc.right_title,
                "right_bullets": sc.right_bullets
            }
        elif sc.slide_type == "chart":
            # Provide default chart data if none exists
            kwargs = {
                "title": sc.title,
                "chart_data": sc.chart_data or {
                    "labels": [],
                    "values": [],
                    "chart_type": "bar"
                }
            }
        elif sc.slide_type == "timeline":
            kwargs = {
                "title": sc.title,
                "events": sc.events
            }
        elif sc.slide_type == "diagram":
            kwargs = {
                "title": sc.title,
                "nodes": sc.nodes,
                "diagram_type": sc.diagram_type
            }
        elif sc.slide_type == "closing":
            kwargs = {
                "title": sc.title or "Thank You",
                "contact": sc.contact,
                "message": sc.message
            }
        else:
            # Default fallback kwargs
            kwargs = {
                "title": sc.title,
                "bullets": sc.bullets
            }

        return kwargs
