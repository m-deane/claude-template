from __future__ import annotations
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

from pptx_generator.config.colors import ColorPalette
from pptx_generator.config.typography import Typography
from pptx_generator.themes.base import BaseTheme
from pptx_generator.slides.base import BaseSlideBuilder


class ChartSlideBuilder(BaseSlideBuilder):
    """Builder for slides with charts."""

    @property
    def slide_type(self) -> str:
        return "chart"

    def _build_content(self, slide: Slide, theme: BaseTheme, palette: ColorPalette,
                       typography: Typography, **kwargs) -> None:
        """Build chart slide."""
        title = kwargs.get('title', 'Chart')
        chart_data_dict = kwargs.get('chart_data', {})
        labels = chart_data_dict.get('labels', ['A', 'B', 'C'])
        values = chart_data_dict.get('values', [10, 20, 15])
        chart_type = chart_data_dict.get('chart_type', 'bar')

        # Title at top
        self._add_text_box(
            slide,
            left=Inches(0.5),
            top=Inches(0.5),
            width=Inches(12.333),
            height=Inches(0.7),
            text=title,
            font_spec=typography.heading,
            alignment=PP_ALIGN.LEFT
        )

        # Accent underline
        line = slide.shapes.add_shape(
            1,  # Line
            Inches(0.5),
            Inches(1.3),
            Inches(2),
            Inches(0)
        )
        line.line.color.rgb = palette.to_rgb(palette.accent)
        line.line.width = Pt(3)

        # Prepare chart data
        chart_data = CategoryChartData()
        chart_data.categories = labels
        chart_data.add_series('Series 1', values)

        # Map chart type to XL_CHART_TYPE
        chart_type_map = {
            'bar': XL_CHART_TYPE.BAR_CLUSTERED,
            'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
            'line': XL_CHART_TYPE.LINE,
            'pie': XL_CHART_TYPE.PIE
        }
        xl_chart_type = chart_type_map.get(chart_type, XL_CHART_TYPE.BAR_CLUSTERED)

        # Add chart
        chart_left = Inches(1)
        chart_top = Inches(2)
        chart_width = Inches(11.333)
        chart_height = Inches(4.5)

        chart_shape = slide.shapes.add_chart(
            xl_chart_type,
            chart_left,
            chart_top,
            chart_width,
            chart_height,
            chart_data
        )

        # Style the chart
        chart = chart_shape.chart

        # Apply color palette to series
        if chart.series:
            series = chart.series[0]
            # Set fill color for bars/columns
            if xl_chart_type in (XL_CHART_TYPE.BAR_CLUSTERED, XL_CHART_TYPE.COLUMN_CLUSTERED):
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = palette.to_rgb(palette.primary)
            elif xl_chart_type == XL_CHART_TYPE.LINE:
                series.format.line.color.rgb = palette.to_rgb(palette.primary)
                series.format.line.width = Pt(2.5)

        # Style chart elements
        if hasattr(chart, 'has_legend'):
            chart.has_legend = False

        # Font styling for chart text
        if hasattr(chart, 'font'):
            chart.font.size = Pt(typography.body.size_pt)
            chart.font.name = typography.body.family
