from __future__ import annotations
import math
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from pptx_generator.config.colors import ColorPalette
from pptx_generator.config.typography import Typography
from pptx_generator.themes.base import BaseTheme
from pptx_generator.slides.base import BaseSlideBuilder


class DiagramSlideBuilder(BaseSlideBuilder):
    """Builder for diagram slides."""

    @property
    def slide_type(self) -> str:
        return "diagram"

    def _build_content(self, slide: Slide, theme: BaseTheme, palette: ColorPalette,
                       typography: Typography, **kwargs) -> None:
        """Build diagram slide."""
        title = kwargs.get('title', 'Diagram')
        nodes = kwargs.get('nodes', [])
        diagram_type = kwargs.get('diagram_type', 'flow')

        # Title at top
        self._add_text_box(
            slide,
            left=Inches(0.5),
            top=Inches(0.5),
            width=Inches(12.333),
            height=Inches(0.7),
            text=title,
            font_spec=typography.heading,
            alignment=PP_ALIGN.LEFT
        )

        # Accent underline
        line = slide.shapes.add_shape(
            1,  # Line
            Inches(0.5),
            Inches(1.3),
            Inches(2),
            Inches(0)
        )
        line.line.color.rgb = palette.to_rgb(palette.accent)
        line.line.width = Pt(3)

        if not nodes:
            return

        # Dispatch to diagram type
        if diagram_type == 'flow':
            self._build_flow_diagram(slide, nodes, palette, typography)
        elif diagram_type == 'hierarchy':
            self._build_hierarchy_diagram(slide, nodes, palette, typography)
        elif diagram_type == 'cycle':
            self._build_cycle_diagram(slide, nodes, palette, typography)

    def _build_flow_diagram(self, slide: Slide, nodes: list[str],
                           palette: ColorPalette, typography: Typography) -> None:
        """Build horizontal flow diagram."""
        num_nodes = len(nodes)
        diagram_top = Inches(2.5)
        diagram_left = Inches(1)
        diagram_width = Inches(11.333)
        node_width = Inches(2)
        node_height = Inches(1)

        spacing = (diagram_width - num_nodes * node_width) / max(1, num_nodes - 1) if num_nodes > 1 else 0

        for idx, node_text in enumerate(nodes):
            node_left = diagram_left + idx * (node_width + spacing)
            node_top = diagram_top

            # Rounded rectangle for node
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                node_left,
                node_top,
                node_width,
                node_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = palette.to_rgb(palette.primary)
            shape.line.color.rgb = palette.to_rgb(palette.accent)
            shape.line.width = Pt(2)

            # Node text
            tf = shape.text_frame
            tf.text = node_text
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.name = typography.body.family
                run.font.size = Pt(typography.body.size_pt)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)

            # Arrow to next node
            if idx < num_nodes - 1:
                arrow_left = node_left + node_width
                arrow_top = node_top + node_height / 2
                arrow_width = spacing
                arrow = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW,
                    arrow_left,
                    arrow_top - Inches(0.15),
                    arrow_width,
                    Inches(0.3)
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = palette.to_rgb(palette.accent)
                arrow.line.fill.background()

    def _build_hierarchy_diagram(self, slide: Slide, nodes: list[str],
                                 palette: ColorPalette, typography: Typography) -> None:
        """Build top-down hierarchy diagram."""
        if not nodes:
            return

        # Simple 2-level hierarchy: first node is parent, rest are children
        parent = nodes[0]
        children = nodes[1:]

        parent_width = Inches(2.5)
        parent_height = Inches(1)
        child_width = Inches(2)
        child_height = Inches(0.8)

        parent_left = Inches(6.667) - parent_width / 2
        parent_top = Inches(2)

        # Parent node
        parent_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            parent_left,
            parent_top,
            parent_width,
            parent_height
        )
        parent_shape.fill.solid()
        parent_shape.fill.fore_color.rgb = palette.to_rgb(palette.primary)
        parent_shape.line.color.rgb = palette.to_rgb(palette.accent)
        parent_shape.line.width = Pt(2)

        tf = parent_shape.text_frame
        tf.text = parent
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.name = typography.body.family
            run.font.size = Pt(typography.body.size_pt)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)

        # Children nodes
        if children:
            num_children = len(children)
            child_y = Inches(4)
            total_width = num_children * child_width + (num_children - 1) * Inches(0.5)
            start_x = Inches(6.667) - total_width / 2

            parent_center_x = parent_left + parent_width / 2
            parent_bottom = parent_top + parent_height

            for idx, child_text in enumerate(children):
                child_left = start_x + idx * (child_width + Inches(0.5))
                child_center_x = child_left + child_width / 2

                # Child node
                child_shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    child_left,
                    child_y,
                    child_width,
                    child_height
                )
                child_shape.fill.solid()
                child_shape.fill.fore_color.rgb = palette.to_rgb(palette.secondary)
                child_shape.line.color.rgb = palette.to_rgb(palette.accent)
                child_shape.line.width = Pt(1.5)

                tf = child_shape.text_frame
                tf.text = child_text
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.name = typography.caption.family
                    run.font.size = Pt(typography.caption.size_pt)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(255, 255, 255)

                # Connecting line
                connector = slide.shapes.add_connector(
                    1,  # Straight connector
                    parent_center_x,
                    parent_bottom,
                    child_center_x,
                    child_y
                )
                connector.line.color.rgb = palette.to_rgb(palette.neutral)
                connector.line.width = Pt(1.5)

    def _build_cycle_diagram(self, slide: Slide, nodes: list[str],
                            palette: ColorPalette, typography: Typography) -> None:
        """Build circular cycle diagram."""
        num_nodes = len(nodes)
        if num_nodes == 0:
            return

        # Circle layout parameters
        center_x = Inches(6.667)
        center_y = Inches(4)
        radius = Inches(2.5)
        node_width = Inches(1.8)
        node_height = Inches(0.8)

        # Calculate positions on circle
        angle_step = 2 * math.pi / num_nodes
        start_angle = -math.pi / 2  # Start at top

        node_positions = []
        for idx in range(num_nodes):
            angle = start_angle + idx * angle_step
            x = center_x + radius * math.cos(angle) - node_width / 2
            y = center_y + radius * math.sin(angle) - node_height / 2
            node_positions.append((x, y, angle))

        # Draw nodes and arrows
        for idx, (node_text, (x, y, angle)) in enumerate(zip(nodes, node_positions)):
            # Node
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x,
                y,
                node_width,
                node_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = palette.to_rgb(palette.primary)
            shape.line.color.rgb = palette.to_rgb(palette.accent)
            shape.line.width = Pt(2)

            tf = shape.text_frame
            tf.text = node_text
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.name = typography.caption.family
                run.font.size = Pt(typography.caption.size_pt - 2)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)

            # Arrow to next node
            next_idx = (idx + 1) % num_nodes
            next_x, next_y, next_angle = node_positions[next_idx]

            # Arrow between nodes (simplified curved arrow)
            mid_angle = angle + angle_step / 2
            arrow_radius = radius * 0.85
            arrow_x = center_x + arrow_radius * math.cos(mid_angle)
            arrow_y = center_y + arrow_radius * math.sin(mid_angle)

            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                arrow_x - Inches(0.3),
                arrow_y - Inches(0.15),
                Inches(0.6),
                Inches(0.3)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = palette.to_rgb(palette.accent)
            arrow.line.fill.background()
            arrow.rotation = math.degrees(mid_angle)
