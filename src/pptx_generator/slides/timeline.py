from __future__ import annotations
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from pptx_generator.config.colors import ColorPalette
from pptx_generator.config.typography import Typography
from pptx_generator.themes.base import BaseTheme
from pptx_generator.slides.base import BaseSlideBuilder


class TimelineSlideBuilder(BaseSlideBuilder):
    """Builder for timeline slides."""

    @property
    def slide_type(self) -> str:
        return "timeline"

    def _build_content(self, slide: Slide, theme: BaseTheme, palette: ColorPalette,
                       typography: Typography, **kwargs) -> None:
        """Build timeline slide."""
        title = kwargs.get('title', 'Timeline')
        events = kwargs.get('events', [])

        # Title at top
        self._add_text_box(
            slide,
            left=Inches(0.5),
            top=Inches(0.5),
            width=Inches(12.333),
            height=Inches(0.7),
            text=title,
            font_spec=typography.heading,
            alignment=PP_ALIGN.LEFT
        )

        # Accent underline
        line = slide.shapes.add_shape(
            1,  # Line
            Inches(0.5),
            Inches(1.3),
            Inches(2),
            Inches(0)
        )
        line.line.color.rgb = palette.to_rgb(palette.accent)
        line.line.width = Pt(3)

        if not events:
            return

        # Timeline bar position
        timeline_y = Inches(3.75)
        timeline_left = Inches(1)
        timeline_right = Inches(12.333)
        timeline_width = timeline_right - timeline_left

        # Horizontal timeline bar
        bar = slide.shapes.add_shape(
            1,  # Line
            timeline_left,
            timeline_y,
            timeline_width,
            Inches(0)
        )
        bar.line.color.rgb = palette.to_rgb(palette.accent)
        bar.line.width = Pt(4)

        # Calculate spacing for events
        num_events = len(events)
        if num_events == 1:
            positions = [timeline_left + timeline_width / 2]
        else:
            spacing = timeline_width / (num_events - 1)
            positions = [timeline_left + i * spacing for i in range(num_events)]

        # Add events
        circle_radius = Inches(0.25)
        label_top = Inches(2)
        desc_top = Inches(4.3)

        for idx, (event, pos) in enumerate(zip(events, positions)):
            label_text = event.get('label', f'Event {idx + 1}')
            description = event.get('description', '')

            # Circle on timeline
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                pos - circle_radius,
                timeline_y - circle_radius,
                circle_radius * 2,
                circle_radius * 2
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = palette.to_rgb(palette.primary)
            circle.line.color.rgb = palette.to_rgb(palette.accent)
            circle.line.width = Pt(2)

            # Label above timeline
            label_width = Inches(2)
            self._add_text_box(
                slide,
                left=pos - label_width / 2,
                top=label_top,
                width=label_width,
                height=Inches(0.5),
                text=label_text,
                font_spec=typography.subheading,
                alignment=PP_ALIGN.CENTER
            )

            # Description below timeline
            desc_width = Inches(2)
            self._add_text_box(
                slide,
                left=pos - desc_width / 2,
                top=desc_top,
                width=desc_width,
                height=Inches(1.5),
                text=description,
                font_spec=typography.caption,
                color_hex=palette.text_light,
                alignment=PP_ALIGN.CENTER
            )
