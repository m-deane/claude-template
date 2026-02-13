from __future__ import annotations
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from pptx_generator.config.colors import ColorPalette
from pptx_generator.config.typography import Typography
from pptx_generator.themes.base import BaseTheme
from pptx_generator.slides.base import BaseSlideBuilder


class AgendaSlideBuilder(BaseSlideBuilder):
    """Builder for agenda slides."""

    @property
    def slide_type(self) -> str:
        return "agenda"

    def _build_content(self, slide: Slide, theme: BaseTheme, palette: ColorPalette,
                       typography: Typography, **kwargs) -> None:
        """Build agenda slide content."""
        items = kwargs.get('items', [])
        title = kwargs.get('title', 'Agenda')

        # Title at top
        title_box = self._add_text_box(
            slide,
            left=Inches(0.5),
            top=Inches(0.5),
            width=Inches(12.333),
            height=Inches(0.8),
            text=title,
            font_spec=typography.heading,
            alignment=PP_ALIGN.LEFT
        )

        # Accent underline for title
        line = slide.shapes.add_shape(
            1,  # Line
            Inches(0.5),
            Inches(1.4),
            Inches(2),
            Inches(0)
        )
        line.line.color.rgb = palette.to_rgb(palette.accent)
        line.line.width = Pt(3)

        # Agenda items with numbered bullets
        start_top = Inches(2)
        item_height = Inches(0.6)
        number_width = Inches(0.8)
        text_left = Inches(1.5)

        for idx, item_text in enumerate(items, 1):
            item_top = start_top + (idx - 1) * item_height

            # Accent-colored number circle
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(0.6),
                item_top + Inches(0.05),
                Inches(0.4),
                Inches(0.4)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = palette.to_rgb(palette.accent)
            circle.line.fill.background()

            # Number text
            tf = circle.text_frame
            tf.text = str(idx)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0]
            run.font.name = typography.body.family
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)

            # Item text
            self._add_text_box(
                slide,
                left=text_left,
                top=item_top,
                width=Inches(10),
                height=item_height,
                text=item_text,
                font_spec=typography.body,
                alignment=PP_ALIGN.LEFT
            )
