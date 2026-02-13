from __future__ import annotations
from abc import ABC, abstractmethod
from pptx.presentation import Presentation
from pptx.slide import Slide
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from pptx_generator.config.colors import ColorPalette
from pptx_generator.config.typography import Typography, FontSpec
from pptx_generator.themes.base import BaseTheme


class BaseSlideBuilder(ABC):
    """Base class for slide type builders."""

    @property
    @abstractmethod
    def slide_type(self) -> str:
        """Return the slide type identifier."""

    def build(self, prs: Presentation, theme: BaseTheme, palette: ColorPalette,
              typography: Typography, slide_number: int | None = None, **kwargs) -> Slide:
        """Build a slide and return it."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        theme.apply_slide_background(slide, palette)
        self._build_content(slide, theme, palette, typography, **kwargs)
        if slide_number is not None:
            theme.add_footer(slide, palette, typography, "", slide_number)
        return slide

    @abstractmethod
    def _build_content(self, slide: Slide, theme: BaseTheme, palette: ColorPalette,
                       typography: Typography, **kwargs) -> None:
        """Build the slide-specific content. Override in subclasses."""

    @staticmethod
    def _add_text_box(slide, left, top, width, height, text, font_spec: FontSpec,
                      color_hex: str | None = None, alignment=PP_ALIGN.LEFT):
        """Helper to add a text box with styled text."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = alignment
        run = p.runs[0] if p.runs else p.add_run()
        if not p.runs[0].text:
            run.text = text
        run.font.name = font_spec.family
        run.font.size = Pt(font_spec.size_pt)
        run.font.bold = font_spec.bold
        run.font.italic = font_spec.italic
        color = color_hex or font_spec.color_hex
        if color:
            hex_clean = color.lstrip('#')
            run.font.color.rgb = RGBColor(int(hex_clean[0:2], 16), int(hex_clean[2:4], 16), int(hex_clean[4:6], 16))
        return txBox

    @staticmethod
    def _add_bullet_list(slide, left, top, width, height, bullets: list[str],
                         font_spec: FontSpec, color_hex: str | None = None,
                         bullet_color_hex: str | None = None):
        """Helper to add a bulleted list."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        color = color_hex or font_spec.color_hex

        for i, bullet_text in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet_text
            p.space_after = Pt(8)
            p.level = 0
            # Set bullet character
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
            buNone = pPr.find(qn('a:buNone'))
            if buNone is not None:
                pPr.remove(buNone)
            buChar = pPr.makeelement(qn('a:buChar'), {'char': '\u2022'})
            pPr.append(buChar)
            if bullet_color_hex:
                hex_clean = bullet_color_hex.lstrip('#')
                buClr = pPr.makeelement(qn('a:buClr'), {})
                srgbClr = buClr.makeelement(qn('a:srgbClr'), {'val': hex_clean})
                buClr.append(srgbClr)
                pPr.append(buClr)

            for run in p.runs:
                run.font.name = font_spec.family
                run.font.size = Pt(font_spec.size_pt)
                run.font.bold = font_spec.bold
                if color:
                    hex_clean = color.lstrip('#')
                    run.font.color.rgb = RGBColor(int(hex_clean[0:2], 16), int(hex_clean[2:4], 16), int(hex_clean[4:6], 16))
        return txBox
