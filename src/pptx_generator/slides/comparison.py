from __future__ import annotations
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from pptx_generator.config.colors import ColorPalette
from pptx_generator.config.typography import Typography
from pptx_generator.themes.base import BaseTheme
from pptx_generator.slides.base import BaseSlideBuilder


class ComparisonSlideBuilder(BaseSlideBuilder):
    """Builder for comparison slides with two columns."""

    @property
    def slide_type(self) -> str:
        return "comparison"

    def _build_content(self, slide: Slide, theme: BaseTheme, palette: ColorPalette,
                       typography: Typography, **kwargs) -> None:
        """Build comparison slide with two columns."""
        title = kwargs.get('title', 'Comparison')
        left_title = kwargs.get('left_title', 'Option A')
        left_bullets = kwargs.get('left_bullets', [])
        right_title = kwargs.get('right_title', 'Option B')
        right_bullets = kwargs.get('right_bullets', [])

        # Main title at top
        self._add_text_box(
            slide,
            left=Inches(0.5),
            top=Inches(0.5),
            width=Inches(12.333),
            height=Inches(0.7),
            text=title,
            font_spec=typography.heading,
            alignment=PP_ALIGN.CENTER
        )

        # Accent underline
        line = slide.shapes.add_shape(
            1,  # Line
            Inches(5.5),
            Inches(1.3),
            Inches(2.333),
            Inches(0)
        )
        line.line.color.rgb = palette.to_rgb(palette.accent)
        line.line.width = Pt(3)

        # Column dimensions
        col_width = Inches(5.5)
        col_top = Inches(1.8)
        left_col_left = Inches(0.5)
        right_col_left = Inches(7.333)
        divider_x = Inches(6.667)

        # Vertical divider line
        divider = slide.shapes.add_shape(
            1,  # Line
            divider_x,
            col_top,
            Inches(0),
            Inches(5)
        )
        divider.line.color.rgb = palette.to_rgb(palette.neutral)
        divider.line.width = Pt(1.5)

        # Left column title
        self._add_text_box(
            slide,
            left=left_col_left,
            top=col_top,
            width=col_width,
            height=Inches(0.5),
            text=left_title,
            font_spec=typography.subheading,
            color_hex=palette.primary,
            alignment=PP_ALIGN.LEFT
        )

        # Left column bullets
        if left_bullets:
            self._add_bullet_list(
                slide,
                left=left_col_left + Inches(0.2),
                top=col_top + Inches(0.6),
                width=col_width - Inches(0.4),
                height=Inches(4.2),
                bullets=left_bullets,
                font_spec=typography.body,
                bullet_color_hex=palette.primary
            )

        # Right column title
        self._add_text_box(
            slide,
            left=right_col_left,
            top=col_top,
            width=col_width,
            height=Inches(0.5),
            text=right_title,
            font_spec=typography.subheading,
            color_hex=palette.secondary,
            alignment=PP_ALIGN.LEFT
        )

        # Right column bullets
        if right_bullets:
            self._add_bullet_list(
                slide,
                left=right_col_left + Inches(0.2),
                top=col_top + Inches(0.6),
                width=col_width - Inches(0.4),
                height=Inches(4.2),
                bullets=right_bullets,
                font_spec=typography.body,
                bullet_color_hex=palette.secondary
            )
