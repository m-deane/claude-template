from __future__ import annotations

from pptx.slide import Slide
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx_generator.config.colors import ColorPalette


def add_stat_card(
    slide: Slide,
    left: int,
    top: int,
    width: int,
    height: int,
    number: str,
    label: str,
    palette: ColorPalette,
    icon_char: str = "",
) -> None:
    """Add a stat card with large number and label.

    Args:
        slide: The slide to add the card to
        left: Left position in EMUs
        top: Top position in EMUs
        width: Width in EMUs
        height: Height in EMUs
        number: The main statistic value
        label: Description label below the number
        palette: ColorPalette for styling
        icon_char: Optional icon character above the number
    """
    # Background card
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = palette.to_rgb(palette.background)
    card.line.color.rgb = palette.to_rgb(palette.neutral)
    card.line.width = Pt(1)

    # Calculate vertical layout
    current_y = top + Emu(228600)  # Start with padding

    # Add icon if provided
    if icon_char:
        icon_size = Emu(685800)  # ~0.75 inches
        icon_x = left + (width - icon_size) // 2

        icon_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, icon_x, current_y, icon_size, icon_size
        )
        icon_circle.fill.solid()
        icon_circle.fill.fore_color.rgb = palette.to_rgb(palette.accent)
        icon_circle.line.color.rgb = palette.to_rgb(palette.accent)

        # Icon text
        icon_text = icon_circle.text_frame
        icon_text.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = icon_text.paragraphs[0]
        p.text = icon_char
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = palette.to_rgb(palette.text_light)
        p.alignment = PP_ALIGN.CENTER

        current_y += icon_size + Emu(228600)

    # Add number
    number_box = slide.shapes.add_textbox(
        left + Emu(228600),
        current_y,
        width - Emu(457200),
        Emu(685800),
    )
    number_frame = number_box.text_frame
    number_frame.word_wrap = False
    p = number_frame.paragraphs[0]
    p.text = number
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = palette.to_rgb(palette.accent)
    p.alignment = PP_ALIGN.CENTER

    current_y += Emu(685800) + Emu(114300)

    # Add label
    label_box = slide.shapes.add_textbox(
        left + Emu(228600),
        current_y,
        width - Emu(457200),
        Emu(457200),
    )
    label_frame = label_box.text_frame
    label_frame.word_wrap = True
    p = label_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(14)
    p.font.color.rgb = palette.to_rgb(palette.text_dark)
    p.alignment = PP_ALIGN.CENTER


def add_stat_row(
    slide: Slide,
    left: int,
    top: int,
    width: int,
    height: int,
    stats: list[dict],
    palette: ColorPalette,
) -> None:
    """Add a row of stat cards.

    Args:
        slide: The slide to add the row to
        left: Left position in EMUs
        top: Top position in EMUs
        width: Width in EMUs
        height: Height in EMUs
        stats: List of dicts with keys 'number', 'label', 'icon_char' (optional)
        palette: ColorPalette for styling
    """
    if not stats:
        return

    num_stats = len(stats)
    spacing = Emu(228600)  # ~0.25 inch spacing
    total_spacing = spacing * (num_stats - 1)
    card_width = (width - total_spacing) // num_stats

    current_x = left

    for stat in stats:
        add_stat_card(
            slide,
            current_x,
            top,
            card_width,
            height,
            stat['number'],
            stat['label'],
            palette,
            stat.get('icon_char', ''),
        )
        current_x += card_width + spacing


def add_progress_bar(
    slide: Slide,
    left: int,
    top: int,
    width: int,
    height: int,
    label: str,
    value: float,
    max_value: float,
    palette: ColorPalette,
) -> None:
    """Add a labeled progress bar.

    Args:
        slide: The slide to add the progress bar to
        left: Left position in EMUs
        top: Top position in EMUs
        width: Width in EMUs
        height: Height in EMUs
        label: Label text on the left
        value: Current value
        max_value: Maximum value
        palette: ColorPalette for styling
    """
    # Add label
    label_width = width // 3
    label_box = slide.shapes.add_textbox(left, top, label_width, height)
    label_frame = label_box.text_frame
    label_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = label_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = palette.to_rgb(palette.text_dark)
    p.alignment = PP_ALIGN.LEFT

    # Bar dimensions
    bar_left = left + label_width + Emu(228600)
    bar_width = width - label_width - Emu(914400)  # Leave space for percentage
    bar_height = height // 2
    bar_top = top + (height - bar_height) // 2

    # Background bar
    bg_bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, bar_left, bar_top, bar_width, bar_height
    )
    bg_bar.fill.solid()
    bg_bar.fill.fore_color.rgb = palette.to_rgb(palette.neutral)
    bg_bar.line.width = Pt(0)

    # Fill bar
    percentage = min(value / max_value, 1.0) if max_value > 0 else 0
    fill_width = int(bar_width * percentage)

    if fill_width > 0:
        fill_bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, bar_left, bar_top, fill_width, bar_height
        )
        fill_bar.fill.solid()
        fill_bar.fill.fore_color.rgb = palette.to_rgb(palette.primary)
        fill_bar.line.width = Pt(0)

    # Percentage label
    pct_left = bar_left + bar_width + Emu(114300)
    pct_width = Emu(685800)
    pct_box = slide.shapes.add_textbox(pct_left, top, pct_width, height)
    pct_frame = pct_box.text_frame
    pct_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = pct_frame.paragraphs[0]
    p.text = f"{int(percentage * 100)}%"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = palette.to_rgb(palette.accent)
    p.alignment = PP_ALIGN.LEFT


def add_icon_text_pair(
    slide: Slide,
    left: int,
    top: int,
    width: int,
    height: int,
    icon_char: str,
    text: str,
    palette: ColorPalette,
) -> None:
    """Add an icon circle paired with text.

    Args:
        slide: The slide to add the pair to
        left: Left position in EMUs
        top: Top position in EMUs
        width: Width in EMUs
        height: Height in EMUs
        icon_char: Icon character to display
        text: Text content
        palette: ColorPalette for styling
    """
    # Icon circle
    icon_size = min(height, Emu(914400))  # Max 1 inch
    icon_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top + (height - icon_size) // 2, icon_size, icon_size
    )
    icon_circle.fill.solid()
    icon_circle.fill.fore_color.rgb = palette.to_rgb(palette.accent)
    icon_circle.line.color.rgb = palette.to_rgb(palette.accent)

    # Icon text
    icon_text = icon_circle.text_frame
    icon_text.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = icon_text.paragraphs[0]
    p.text = icon_char
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = palette.to_rgb(palette.text_light)
    p.alignment = PP_ALIGN.CENTER

    # Text content
    text_left = left + icon_size + Emu(228600)
    text_width = width - icon_size - Emu(228600)
    text_box = slide.shapes.add_textbox(text_left, top, text_width, height)
    text_frame = text_box.text_frame
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.color.rgb = palette.to_rgb(palette.text_dark)
    p.alignment = PP_ALIGN.LEFT


def add_kpi_dashboard(
    slide: Slide,
    left: int,
    top: int,
    width: int,
    height: int,
    kpis: list[dict],
    palette: ColorPalette,
) -> None:
    """Add a KPI dashboard grid layout.

    Args:
        slide: The slide to add the dashboard to
        left: Left position in EMUs
        top: Top position in EMUs
        width: Width in EMUs
        height: Height in EMUs
        kpis: List of dicts with keys 'title', 'value', 'change', 'trend'
              trend values: 'up', 'down', 'flat'
        palette: ColorPalette for styling
    """
    if not kpis:
        return

    # Calculate grid layout (2 columns)
    num_kpis = len(kpis)
    cols = 2
    rows = (num_kpis + cols - 1) // cols

    spacing = Emu(228600)
    card_width = (width - spacing) // cols
    card_height = (height - spacing * (rows - 1)) // rows

    trend_colors = {
        'up': palette.success,
        'down': palette.danger,
        'flat': palette.neutral,
    }

    trend_symbols = {
        'up': '▲',
        'down': '▼',
        'flat': '▬',
    }

    for i, kpi in enumerate(kpis):
        row = i // cols
        col = i % cols

        card_left = left + col * (card_width + spacing)
        card_top = top + row * (card_height + spacing)

        # Background card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            card_left,
            card_top,
            card_width,
            card_height,
        )
        card.fill.solid()
        card.fill.fore_color.rgb = palette.to_rgb(palette.background)
        card.line.color.rgb = palette.to_rgb(palette.neutral)
        card.line.width = Pt(1.5)

        # Title
        title_box = slide.shapes.add_textbox(
            card_left + Emu(228600),
            card_top + Emu(228600),
            card_width - Emu(457200),
            Emu(457200),
        )
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = kpi['title']
        p.font.size = Pt(12)
        p.font.color.rgb = palette.to_rgb(palette.text_dark)
        p.alignment = PP_ALIGN.LEFT

        # Value
        value_box = slide.shapes.add_textbox(
            card_left + Emu(228600),
            card_top + Emu(685800),
            card_width - Emu(457200),
            Emu(685800),
        )
        value_frame = value_box.text_frame
        p = value_frame.paragraphs[0]
        p.text = kpi['value']
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = palette.to_rgb(palette.primary)
        p.alignment = PP_ALIGN.LEFT

        # Change with trend
        trend = kpi.get('trend', 'flat')
        change_text = f"{trend_symbols.get(trend, '')} {kpi['change']}"

        change_box = slide.shapes.add_textbox(
            card_left + Emu(228600),
            card_top + card_height - Emu(571500),
            card_width - Emu(457200),
            Emu(343200),
        )
        change_frame = change_box.text_frame
        p = change_frame.paragraphs[0]
        p.text = change_text
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = palette.to_rgb(trend_colors.get(trend, palette.neutral))
        p.alignment = PP_ALIGN.LEFT
