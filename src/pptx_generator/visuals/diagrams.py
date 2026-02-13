from __future__ import annotations

import math
from pptx.slide import Slide
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx_generator.config.colors import ColorPalette


def add_flow_diagram(
    slide: Slide,
    left: int,
    top: int,
    width: int,
    height: int,
    nodes: list[str],
    palette: ColorPalette,
    node_style: str = "rounded",
) -> None:
    """Add a horizontal flow diagram with connected nodes.

    Args:
        slide: The slide to add the diagram to
        left: Left position in EMUs
        top: Top position in EMUs
        width: Width in EMUs
        height: Height in EMUs
        nodes: List of node labels
        palette: ColorPalette for styling
        node_style: Shape style - 'rounded' or 'rectangle'
    """
    if not nodes:
        return

    num_nodes = len(nodes)
    arrow_width = Emu(914400)  # ~1 inch for arrows
    spacing = Emu(457200)  # ~0.5 inch between elements

    # Calculate node width
    total_arrow_space = arrow_width * (num_nodes - 1) if num_nodes > 1 else 0
    total_spacing = spacing * (num_nodes - 1) if num_nodes > 1 else 0
    available_width = width - total_arrow_space - total_spacing
    node_width = available_width // num_nodes
    node_height = height

    # Choose shape type
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if node_style == "rounded" else MSO_SHAPE.RECTANGLE

    current_x = left

    for i, label in enumerate(nodes):
        # Add node
        node = slide.shapes.add_shape(
            shape_type, current_x, top, node_width, node_height
        )

        # Style node
        node.fill.solid()
        node.fill.fore_color.rgb = palette.to_rgb(palette.primary)
        node.line.color.rgb = palette.to_rgb(palette.primary)
        node.line.width = Pt(1)

        # Add text
        text_frame = node.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = text_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = palette.to_rgb(palette.text_light)
        p.alignment = PP_ALIGN.CENTER

        # Add arrow between nodes
        if i < num_nodes - 1:
            arrow_x = current_x + node_width + spacing
            arrow_y = top + (node_height // 2) - (Emu(457200) // 2)

            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                arrow_x,
                arrow_y,
                arrow_width,
                Emu(457200),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = palette.to_rgb(palette.accent)
            arrow.line.color.rgb = palette.to_rgb(palette.accent)

        current_x += node_width + spacing + arrow_width


def add_hierarchy_diagram(
    slide: Slide,
    left: int,
    top: int,
    width: int,
    height: int,
    root: str,
    children: list[str],
    palette: ColorPalette,
) -> None:
    """Add a top-down hierarchy diagram with root node and children.

    Args:
        slide: The slide to add the diagram to
        left: Left position in EMUs
        top: Top position in EMUs
        width: Width in EMUs
        height: Height in EMUs
        root: Root node label
        children: List of child node labels
        palette: ColorPalette for styling
    """
    if not children:
        return

    # Root node dimensions (larger)
    root_width = width // 3
    root_height = height // 3
    root_x = left + (width - root_width) // 2
    root_y = top

    # Add root node
    root_node = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, root_x, root_y, root_width, root_height
    )
    root_node.fill.solid()
    root_node.fill.fore_color.rgb = palette.to_rgb(palette.primary)
    root_node.line.color.rgb = palette.to_rgb(palette.primary)
    root_node.line.width = Pt(2)

    # Root text
    text_frame = root_node.text_frame
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = text_frame.paragraphs[0]
    p.text = root
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = palette.to_rgb(palette.text_light)
    p.alignment = PP_ALIGN.CENTER

    # Child nodes
    num_children = len(children)
    child_width = (width - Emu(914400)) // num_children  # Leave spacing
    child_height = height // 3
    child_y = top + height - child_height

    spacing = (width - (child_width * num_children)) // (num_children + 1)

    root_center_x = root_x + root_width // 2
    root_bottom_y = root_y + root_height

    for i, child_label in enumerate(children):
        child_x = left + spacing + (i * (child_width + spacing))

        # Add child node
        child_node = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, child_x, child_y, child_width, child_height
        )
        child_node.fill.solid()
        child_node.fill.fore_color.rgb = palette.to_rgb(palette.secondary)
        child_node.line.color.rgb = palette.to_rgb(palette.secondary)
        child_node.line.width = Pt(1.5)

        # Child text
        text_frame = child_node.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = text_frame.paragraphs[0]
        p.text = child_label
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = palette.to_rgb(palette.text_light)
        p.alignment = PP_ALIGN.CENTER

        # Add connecting line
        child_center_x = child_x + child_width // 2
        child_top_y = child_y

        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            root_center_x,
            root_bottom_y,
            child_center_x,
            child_top_y,
        )
        connector.line.color.rgb = palette.to_rgb(palette.neutral)
        connector.line.width = Pt(2)


def add_cycle_diagram(
    slide: Slide,
    left: int,
    top: int,
    width: int,
    height: int,
    steps: list[str],
    palette: ColorPalette,
) -> None:
    """Add a circular cycle diagram with steps arranged in a circle.

    Args:
        slide: The slide to add the diagram to
        left: Left position in EMUs
        top: Top position in EMUs
        width: Width in EMUs
        height: Height in EMUs
        steps: List of step labels
        palette: ColorPalette for styling
    """
    if not steps:
        return

    num_steps = len(steps)
    center_x = left + width // 2
    center_y = top + height // 2

    # Calculate radius for node placement (use 80% of available space)
    radius = min(width, height) // 2 * 0.7

    # Node dimensions
    node_width = Emu(1143000)  # ~1.25 inches
    node_height = Emu(685800)  # ~0.75 inches

    colors = [palette.primary, palette.accent, palette.secondary, palette.success, palette.warning]

    for i, step_label in enumerate(steps):
        # Calculate position on circle
        angle = (2 * math.pi * i / num_steps) - (math.pi / 2)  # Start at top
        node_x = int(center_x + radius * math.cos(angle) - node_width // 2)
        node_y = int(center_y + radius * math.sin(angle) - node_height // 2)

        # Add node
        node = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, node_x, node_y, node_width, node_height
        )

        # Style node
        color_hex = colors[i % len(colors)]
        node.fill.solid()
        node.fill.fore_color.rgb = palette.to_rgb(color_hex)
        node.line.color.rgb = palette.to_rgb(color_hex)
        node.line.width = Pt(2)

        # Add text
        text_frame = node.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = step_label
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = palette.to_rgb(palette.text_light)
        p.alignment = PP_ALIGN.CENTER

        # Add curved arrow to next step
        if num_steps > 1:
            next_angle = (2 * math.pi * ((i + 1) % num_steps) / num_steps) - (math.pi / 2)
            mid_angle = (angle + next_angle) / 2

            # Use smaller radius for arrow placement
            arrow_radius = radius * 0.85
            arrow_x = int(center_x + arrow_radius * math.cos(mid_angle))
            arrow_y = int(center_y + arrow_radius * math.sin(mid_angle))

            # Add small arrow shape
            arrow_size = Emu(228600)  # ~0.25 inches
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                arrow_x - arrow_size // 2,
                arrow_y - arrow_size // 4,
                arrow_size,
                arrow_size // 2,
            )

            # Rotate arrow to point in direction of flow
            arrow.rotation = math.degrees(mid_angle) + 90

            arrow.fill.solid()
            arrow.fill.fore_color.rgb = palette.to_rgb(palette.accent)
            arrow.line.color.rgb = palette.to_rgb(palette.accent)


def add_process_arrows(
    slide: Slide,
    left: int,
    top: int,
    width: int,
    height: int,
    steps: list[str],
    palette: ColorPalette,
) -> None:
    """Add a chevron/arrow process diagram.

    Args:
        slide: The slide to add the diagram to
        left: Left position in EMUs
        top: Top position in EMUs
        width: Width in EMUs
        height: Height in EMUs
        steps: List of step labels
        palette: ColorPalette for styling
    """
    if not steps:
        return

    num_steps = len(steps)
    chevron_width = width // num_steps
    overlap = Emu(228600)  # ~0.25 inch overlap for chevron effect

    colors = [palette.primary, palette.secondary]

    current_x = left

    for i, step_label in enumerate(steps):
        # Add chevron
        chevron = slide.shapes.add_shape(
            MSO_SHAPE.CHEVRON,
            current_x - (overlap if i > 0 else 0),
            top,
            chevron_width + overlap,
            height,
        )

        # Alternate colors
        color_hex = colors[i % len(colors)]
        chevron.fill.solid()
        chevron.fill.fore_color.rgb = palette.to_rgb(color_hex)
        chevron.line.color.rgb = palette.to_rgb(color_hex)
        chevron.line.width = Pt(1)

        # Add text
        text_frame = chevron.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        text_frame.word_wrap = True
        text_frame.margin_left = Emu(228600)
        text_frame.margin_right = Emu(457200)  # Extra margin on right for chevron point

        p = text_frame.paragraphs[0]
        p.text = step_label
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = palette.to_rgb(palette.text_light)
        p.alignment = PP_ALIGN.CENTER

        current_x += chevron_width
