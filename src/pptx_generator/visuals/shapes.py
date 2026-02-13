from __future__ import annotations

from pptx.slide import Slide
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx_generator.config.colors import ColorPalette


def add_rounded_rectangle(
    slide: Slide,
    left: int,
    top: int,
    width: int,
    height: int,
    fill_hex: str,
    border_hex: str | None = None,
    text: str = "",
    font_size_pt: int = 14,
    font_color_hex: str = "#FFFFFF",
):
    """Add a rounded rectangle shape with optional text.

    Args:
        slide: The slide to add the shape to
        left: Left position in EMUs
        top: Top position in EMUs
        width: Width in EMUs
        height: Height in EMUs
        fill_hex: Fill color as hex string
        border_hex: Optional border color as hex string
        text: Optional text content
        font_size_pt: Font size in points
        font_color_hex: Font color as hex string

    Returns:
        The created shape object
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )

    # Fill
    shape.fill.solid()
    fill_color = fill_hex.lstrip('#')
    shape.fill.fore_color.rgb = RGBColor.from_string(fill_color)

    # Border
    if border_hex:
        border_color = border_hex.lstrip('#')
        shape.line.color.rgb = RGBColor.from_string(border_color)
        shape.line.width = Pt(1.5)
    else:
        shape.line.width = Pt(0)

    # Text
    if text:
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size_pt)
        p.font.bold = True

        font_color = font_color_hex.lstrip('#')
        p.font.color.rgb = RGBColor.from_string(font_color)
        p.alignment = PP_ALIGN.CENTER

    return shape


def add_circle(
    slide: Slide,
    left: int,
    top: int,
    diameter: int,
    fill_hex: str,
    border_hex: str | None = None,
    text: str = "",
    font_size_pt: int = 12,
    font_color_hex: str = "#FFFFFF",
):
    """Add a circle shape with optional text.

    Args:
        slide: The slide to add the shape to
        left: Left position in EMUs
        top: Top position in EMUs
        diameter: Diameter in EMUs
        fill_hex: Fill color as hex string
        border_hex: Optional border color as hex string
        text: Optional text content
        font_size_pt: Font size in points
        font_color_hex: Font color as hex string

    Returns:
        The created shape object
    """
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)

    # Fill
    shape.fill.solid()
    fill_color = fill_hex.lstrip('#')
    shape.fill.fore_color.rgb = RGBColor.from_string(fill_color)

    # Border
    if border_hex:
        border_color = border_hex.lstrip('#')
        shape.line.color.rgb = RGBColor.from_string(border_color)
        shape.line.width = Pt(2)
    else:
        shape.line.width = Pt(0)

    # Text
    if text:
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size_pt)
        p.font.bold = True

        font_color = font_color_hex.lstrip('#')
        p.font.color.rgb = RGBColor.from_string(font_color)
        p.alignment = PP_ALIGN.CENTER

    return shape


def add_arrow(
    slide: Slide,
    left: int,
    top: int,
    width: int,
    height: int,
    fill_hex: str,
    direction: str = "right",
):
    """Add an arrow shape.

    Args:
        slide: The slide to add the shape to
        left: Left position in EMUs
        top: Top position in EMUs
        width: Width in EMUs
        height: Height in EMUs
        fill_hex: Fill color as hex string
        direction: Arrow direction - 'right', 'left', 'up', 'down'

    Returns:
        The created shape object
    """
    # Map direction to shape type
    arrow_shapes = {
        'right': MSO_SHAPE.RIGHT_ARROW,
        'left': MSO_SHAPE.LEFT_ARROW,
        'up': MSO_SHAPE.UP_ARROW,
        'down': MSO_SHAPE.DOWN_ARROW,
    }

    shape_type = arrow_shapes.get(direction.lower(), MSO_SHAPE.RIGHT_ARROW)
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)

    # Fill
    shape.fill.solid()
    fill_color = fill_hex.lstrip('#')
    shape.fill.fore_color.rgb = RGBColor.from_string(fill_color)

    # Border (same color as fill)
    shape.line.color.rgb = RGBColor.from_string(fill_color)
    shape.line.width = Pt(1)

    return shape


def add_line(
    slide: Slide,
    start_x: int,
    start_y: int,
    end_x: int,
    end_y: int,
    color_hex: str,
    width_pt: float = 1.0,
):
    """Add a line connector.

    Args:
        slide: The slide to add the line to
        start_x: Starting X position in EMUs
        start_y: Starting Y position in EMUs
        end_x: Ending X position in EMUs
        end_y: Ending Y position in EMUs
        color_hex: Line color as hex string
        width_pt: Line width in points

    Returns:
        The created connector object
    """
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y
    )

    # Line style
    line_color = color_hex.lstrip('#')
    connector.line.color.rgb = RGBColor.from_string(line_color)
    connector.line.width = Pt(width_pt)

    return connector


def add_chevron(
    slide: Slide,
    left: int,
    top: int,
    width: int,
    height: int,
    fill_hex: str,
    text: str = "",
    font_color_hex: str = "#FFFFFF",
):
    """Add a chevron/pentagon shape with optional text.

    Args:
        slide: The slide to add the shape to
        left: Left position in EMUs
        top: Top position in EMUs
        width: Width in EMUs
        height: Height in EMUs
        fill_hex: Fill color as hex string
        text: Optional text content
        font_color_hex: Font color as hex string

    Returns:
        The created shape object
    """
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)

    # Fill
    shape.fill.solid()
    fill_color = fill_hex.lstrip('#')
    shape.fill.fore_color.rgb = RGBColor.from_string(fill_color)

    # Border (same color as fill)
    shape.line.color.rgb = RGBColor.from_string(fill_color)
    shape.line.width = Pt(1)

    # Text
    if text:
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        text_frame.margin_left = Emu(228600)
        text_frame.margin_right = Emu(457200)

        p = text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(14)
        p.font.bold = True

        font_color = font_color_hex.lstrip('#')
        p.font.color.rgb = RGBColor.from_string(font_color)
        p.alignment = PP_ALIGN.CENTER

    return shape


def add_accent_divider(
    slide: Slide, left: int, top: int, width: int, color_hex: str, thickness_pt: int = 2
):
    """Add a thin horizontal line divider.

    Args:
        slide: The slide to add the divider to
        left: Left position in EMUs
        top: Top position in EMUs (center of the line)
        width: Width in EMUs
        color_hex: Line color as hex string
        thickness_pt: Line thickness in points

    Returns:
        The created shape object
    """
    # Create a thin rectangle as a divider
    height = Pt(thickness_pt)

    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top - height // 2, width, height
    )

    # Fill
    divider.fill.solid()
    line_color = color_hex.lstrip('#')
    divider.fill.fore_color.rgb = RGBColor.from_string(line_color)

    # No border
    divider.line.width = Pt(0)

    return divider


def add_triangle(
    slide: Slide,
    left: int,
    top: int,
    width: int,
    height: int,
    fill_hex: str,
    border_hex: str | None = None,
    direction: str = "up",
):
    """Add a triangle shape.

    Args:
        slide: The slide to add the shape to
        left: Left position in EMUs
        top: Top position in EMUs
        width: Width in EMUs
        height: Height in EMUs
        fill_hex: Fill color as hex string
        border_hex: Optional border color as hex string
        direction: Triangle direction - 'up', 'down', 'left', 'right'

    Returns:
        The created shape object
    """
    # Map direction to shape type
    triangle_shapes = {
        'up': MSO_SHAPE.ISOSCELES_TRIANGLE,
        'down': MSO_SHAPE.ISOSCELES_TRIANGLE,
        'left': MSO_SHAPE.ISOSCELES_TRIANGLE,
        'right': MSO_SHAPE.ISOSCELES_TRIANGLE,
    }

    shape = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, left, top, width, height)

    # Rotate based on direction (default is pointing up)
    rotation_map = {
        'up': 0,
        'right': 90,
        'down': 180,
        'left': 270,
    }
    shape.rotation = rotation_map.get(direction.lower(), 0)

    # Fill
    shape.fill.solid()
    fill_color = fill_hex.lstrip('#')
    shape.fill.fore_color.rgb = RGBColor.from_string(fill_color)

    # Border
    if border_hex:
        border_color = border_hex.lstrip('#')
        shape.line.color.rgb = RGBColor.from_string(border_color)
        shape.line.width = Pt(1.5)
    else:
        shape.line.width = Pt(0)

    return shape


def add_star(
    slide: Slide,
    left: int,
    top: int,
    width: int,
    height: int,
    fill_hex: str,
    border_hex: str | None = None,
    points: int = 5,
):
    """Add a star shape.

    Args:
        slide: The slide to add the shape to
        left: Left position in EMUs
        top: Top position in EMUs
        width: Width in EMUs
        height: Height in EMUs
        fill_hex: Fill color as hex string
        border_hex: Optional border color as hex string
        points: Number of star points (5 or 7)

    Returns:
        The created shape object
    """
    # Use 5-point or 7-point star
    shape_type = MSO_SHAPE.STAR_5 if points == 5 else MSO_SHAPE.STAR_7
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)

    # Fill
    shape.fill.solid()
    fill_color = fill_hex.lstrip('#')
    shape.fill.fore_color.rgb = RGBColor.from_string(fill_color)

    # Border
    if border_hex:
        border_color = border_hex.lstrip('#')
        shape.line.color.rgb = RGBColor.from_string(border_color)
        shape.line.width = Pt(1.5)
    else:
        shape.line.width = Pt(0)

    return shape


def add_hexagon(
    slide: Slide,
    left: int,
    top: int,
    width: int,
    height: int,
    fill_hex: str,
    border_hex: str | None = None,
    text: str = "",
    font_size_pt: int = 14,
    font_color_hex: str = "#FFFFFF",
):
    """Add a hexagon shape with optional text.

    Args:
        slide: The slide to add the shape to
        left: Left position in EMUs
        top: Top position in EMUs
        width: Width in EMUs
        height: Height in EMUs
        fill_hex: Fill color as hex string
        border_hex: Optional border color as hex string
        text: Optional text content
        font_size_pt: Font size in points
        font_color_hex: Font color as hex string

    Returns:
        The created shape object
    """
    shape = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, left, top, width, height)

    # Fill
    shape.fill.solid()
    fill_color = fill_hex.lstrip('#')
    shape.fill.fore_color.rgb = RGBColor.from_string(fill_color)

    # Border
    if border_hex:
        border_color = border_hex.lstrip('#')
        shape.line.color.rgb = RGBColor.from_string(border_color)
        shape.line.width = Pt(1.5)
    else:
        shape.line.width = Pt(0)

    # Text
    if text:
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size_pt)
        p.font.bold = True

        font_color = font_color_hex.lstrip('#')
        p.font.color.rgb = RGBColor.from_string(font_color)
        p.alignment = PP_ALIGN.CENTER

    return shape
