from __future__ import annotations

from pptx.slide import Slide
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.dml import MSO_LINE
from pptx_generator.config.colors import ColorPalette


def add_bar_chart(
    slide: Slide,
    left: int,
    top: int,
    width: int,
    height: int,
    categories: list[str],
    series_data: dict[str, list[float]],
    palette: ColorPalette,
    title: str = "",
) -> None:
    """Add a clustered column/bar chart to the slide.

    Args:
        slide: The slide to add the chart to
        left: Left position in EMUs
        top: Top position in EMUs
        width: Width in EMUs
        height: Height in EMUs
        categories: List of category labels for x-axis
        series_data: Dictionary mapping series names to value lists
        palette: ColorPalette for styling
        title: Optional chart title
    """
    chart_data = CategoryChartData()
    chart_data.categories = categories

    for name, values in series_data.items():
        chart_data.add_series(name, values)

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_shape.chart

    # Configure title
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.paragraphs[0].text = title
        chart.chart_title.text_frame.paragraphs[0].font.bold = True
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(18)
        chart.chart_title.text_frame.paragraphs[0].font.color.rgb = palette.to_rgb(palette.text_dark)

    # Hide legend if single series
    if len(series_data) == 1:
        chart.has_legend = False
    else:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(10)

    # Style series with palette colors
    plot = chart.plots[0]
    series_colors = [palette.primary, palette.secondary, palette.accent, palette.success]

    for i, series in enumerate(plot.series):
        color_hex = series_colors[i % len(series_colors)]
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = palette.to_rgb(color_hex)

        # Add data labels for single series charts
        if len(series_data) == 1:
            series.has_data_labels = True
            data_labels = series.data_labels
            data_labels.font.size = Pt(10)
            data_labels.font.color.rgb = palette.to_rgb(palette.text_dark)

    # Clean up gridlines
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = palette.to_rgb(palette.neutral)
    chart.value_axis.major_gridlines.format.line.width = Pt(0.5)


def add_line_chart(
    slide: Slide,
    left: int,
    top: int,
    width: int,
    height: int,
    categories: list[str],
    series_data: dict[str, list[float]],
    palette: ColorPalette,
    title: str = "",
    show_markers: bool = True,
) -> None:
    """Add a line chart to the slide.

    Args:
        slide: The slide to add the chart to
        left: Left position in EMUs
        top: Top position in EMUs
        width: Width in EMUs
        height: Height in EMUs
        categories: List of category labels for x-axis
        series_data: Dictionary mapping series names to value lists
        palette: ColorPalette for styling
        title: Optional chart title
        show_markers: Whether to show data point markers
    """
    chart_data = CategoryChartData()
    chart_data.categories = categories

    for name, values in series_data.items():
        chart_data.add_series(name, values)

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS if show_markers else XL_CHART_TYPE.LINE,
        left,
        top,
        width,
        height,
        chart_data,
    )
    chart = chart_shape.chart

    # Configure title
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.paragraphs[0].text = title
        chart.chart_title.text_frame.paragraphs[0].font.bold = True
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(18)
        chart.chart_title.text_frame.paragraphs[0].font.color.rgb = palette.to_rgb(palette.text_dark)

    # Configure legend
    if len(series_data) == 1:
        chart.has_legend = False
    else:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(10)

    # Style series with palette colors
    plot = chart.plots[0]
    series_colors = [palette.primary, palette.accent, palette.secondary, palette.success]

    for i, series in enumerate(plot.series):
        color_hex = series_colors[i % len(series_colors)]
        color_rgb = palette.to_rgb(color_hex)

        # Line color
        series.format.line.color.rgb = color_rgb
        series.format.line.width = Pt(2.5)

        # Marker styling
        if show_markers and hasattr(series, 'marker'):
            series.marker.format.fill.solid()
            series.marker.format.fill.fore_color.rgb = color_rgb
            series.marker.format.line.color.rgb = color_rgb
            series.marker.size = 6

    # Clean up gridlines
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = palette.to_rgb(palette.neutral)
    chart.value_axis.major_gridlines.format.line.width = Pt(0.5)


def add_pie_chart(
    slide: Slide,
    left: int,
    top: int,
    width: int,
    height: int,
    categories: list[str],
    values: list[float],
    palette: ColorPalette,
    title: str = "",
    show_percentages: bool = True,
) -> None:
    """Add a pie chart to the slide.

    Args:
        slide: The slide to add the chart to
        left: Left position in EMUs
        top: Top position in EMUs
        width: Width in EMUs
        height: Height in EMUs
        categories: List of category labels
        values: List of values for each category
        palette: ColorPalette for styling
        title: Optional chart title
        show_percentages: Whether to show percentage labels
    """
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('Series 1', values)

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, left, top, width, height, chart_data
    )
    chart = chart_shape.chart

    # Configure title
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.paragraphs[0].text = title
        chart.chart_title.text_frame.paragraphs[0].font.bold = True
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(18)
        chart.chart_title.text_frame.paragraphs[0].font.color.rgb = palette.to_rgb(palette.text_dark)

    # Configure legend
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.font.size = Pt(10)

    # Style pie slices with palette colors
    plot = chart.plots[0]
    slice_colors = [
        palette.primary,
        palette.accent,
        palette.secondary,
        palette.success,
        palette.warning,
        palette.neutral,
    ]

    for i, point in enumerate(plot.series[0].points):
        color_hex = slice_colors[i % len(slice_colors)]
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = palette.to_rgb(color_hex)

    # Add data labels
    if show_percentages:
        plot.series[0].has_data_labels = True
        data_labels = plot.series[0].data_labels
        data_labels.font.size = Pt(10)
        data_labels.font.color.rgb = palette.to_rgb(palette.text_light)
        data_labels.font.bold = True
        data_labels.position = XL_LABEL_POSITION.INSIDE_END


def add_horizontal_bar_chart(
    slide: Slide,
    left: int,
    top: int,
    width: int,
    height: int,
    categories: list[str],
    series_data: dict[str, list[float]],
    palette: ColorPalette,
    title: str = "",
) -> None:
    """Add a horizontal bar chart to the slide.

    Args:
        slide: The slide to add the chart to
        left: Left position in EMUs
        top: Top position in EMUs
        width: Width in EMUs
        height: Height in EMUs
        categories: List of category labels for y-axis
        series_data: Dictionary mapping series names to value lists
        palette: ColorPalette for styling
        title: Optional chart title
    """
    chart_data = CategoryChartData()
    chart_data.categories = categories

    for name, values in series_data.items():
        chart_data.add_series(name, values)

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_shape.chart

    # Configure title
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.paragraphs[0].text = title
        chart.chart_title.text_frame.paragraphs[0].font.bold = True
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(18)
        chart.chart_title.text_frame.paragraphs[0].font.color.rgb = palette.to_rgb(palette.text_dark)

    # Hide legend if single series
    if len(series_data) == 1:
        chart.has_legend = False
    else:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(10)

    # Style series with palette colors
    plot = chart.plots[0]
    series_colors = [palette.primary, palette.secondary, palette.accent, palette.success]

    for i, series in enumerate(plot.series):
        color_hex = series_colors[i % len(series_colors)]
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = palette.to_rgb(color_hex)

        # Add data labels for single series charts
        if len(series_data) == 1:
            series.has_data_labels = True
            data_labels = series.data_labels
            data_labels.font.size = Pt(10)
            data_labels.font.color.rgb = palette.to_rgb(palette.text_dark)

    # Clean up gridlines
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = palette.to_rgb(palette.neutral)
    chart.value_axis.major_gridlines.format.line.width = Pt(0.5)
