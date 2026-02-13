"""Corporate theme - clean, professional styling."""

from __future__ import annotations

from typing import TYPE_CHECKING

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from pptx_generator.themes.base import BaseTheme

if TYPE_CHECKING:
    from pptx.slide import Slide
    from pptx.shapes.base import BaseShape

    from pptx_generator.config.colors import ColorPalette
    from pptx_generator.config.typography import Typography


class CorporateTheme(BaseTheme):
    """Corporate theme with clean, professional styling."""

    @property
    def name(self) -> str:
        """Return the theme name."""
        return "corporate"

    def apply_slide_background(self, slide: Slide, palette: ColorPalette) -> None:
        """Apply white/light background to the slide.

        Args:
            slide: The slide to style
            palette: Color palette configuration
        """
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self._hex_to_rgb(palette.background)

    def add_accent_bar(
        self, slide: Slide, palette: ColorPalette, position: str = "left"
    ) -> None:
        """Add navy accent bar to the slide.

        Args:
            slide: The slide to modify
            palette: Color palette configuration
            position: Position of the accent bar (defaults to "left")
        """
        width, height = self._get_slide_dimensions(slide)

        if position == "left":
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0),
                Inches(0),
                Inches(0.15),
                height,
            )
        elif position == "right":
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                width - Inches(0.15),
                Inches(0),
                Inches(0.15),
                height,
            )
        elif position == "top":
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0),
                Inches(0),
                width,
                Inches(0.15),
            )
        elif position == "bottom":
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0),
                height - Inches(0.15),
                width,
                Inches(0.15),
            )
        else:
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0),
                Inches(0),
                Inches(0.15),
                height,
            )

        self._set_shape_fill(bar, palette.primary)
        bar.line.fill.background()

    def add_footer(
        self,
        slide: Slide,
        palette: ColorPalette,
        typography: Typography,
        text: str,
        slide_number: int | None = None,
    ) -> None:
        """Add footer with thin line separator.

        Args:
            slide: The slide to modify
            palette: Color palette configuration
            typography: Typography configuration
            text: Footer text content
            slide_number: Optional slide number to display
        """
        width, height = self._get_slide_dimensions(slide)

        # Add thin separator line
        separator = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5),
            height - Inches(0.6),
            width - Inches(1.0),
            Pt(1),
        )
        self._set_shape_fill(separator, palette.neutral)
        separator.line.fill.background()

        # Add footer text
        footer_text = text
        if slide_number is not None:
            footer_text = f"{text}  |  Slide {slide_number}"

        footer_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5),
            height - Inches(0.5),
            width - Inches(1.0),
            Inches(0.35),
        )
        footer_box.fill.background()
        footer_box.line.fill.background()

        self._add_text_to_shape(
            footer_box,
            footer_text,
            typography.caption.family,
            typography.caption.size_pt,
            typography.caption.color_hex or palette.text_dark,
            bold=typography.caption.bold,
            italic=typography.caption.italic,
        )

        text_frame = footer_box.text_frame
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def style_shape(
        self, shape: BaseShape, palette: ColorPalette, style: str = "primary"
    ) -> None:
        """Apply corporate styling to a shape with slight transparency.

        Args:
            shape: The shape to style
            palette: Color palette configuration
            style: Style variant ("primary", "secondary", "accent", "neutral")
        """
        color_map = {
            "primary": palette.primary,
            "secondary": palette.secondary,
            "accent": palette.accent,
            "neutral": palette.neutral,
        }

        fill_color = color_map.get(style, palette.primary)
        self._set_shape_fill(shape, fill_color, transparency=0.1)
        self._set_shape_line(shape, fill_color, width_pt=1.5)

    def add_decorative_element(
        self, slide: Slide, palette: ColorPalette, element_type: str = "corner"
    ) -> None:
        """Add small colored triangle in bottom-right corner.

        Args:
            slide: The slide to modify
            palette: Color palette configuration
            element_type: Type of decoration (defaults to "corner")
        """
        if element_type == "corner":
            width, height = self._get_slide_dimensions(slide)

            triangle = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_TRIANGLE,
                width - Inches(1.2),
                height - Inches(1.0),
                Inches(0.8),
                Inches(0.6),
            )
            self._set_shape_fill(triangle, palette.accent, transparency=0.3)
            triangle.line.fill.background()

            # Rotate to point upward-left
            triangle.rotation = 45
        elif element_type == "none":
            pass
        else:
            # Default to corner decoration
            self.add_decorative_element(slide, palette, "corner")
