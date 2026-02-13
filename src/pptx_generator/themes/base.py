"""Abstract base theme for PowerPoint presentations."""

from __future__ import annotations

from abc import ABC, abstractmethod
from typing import TYPE_CHECKING

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

if TYPE_CHECKING:
    from pptx.slide import Slide
    from pptx.shapes.base import BaseShape

    from pptx_generator.config.colors import ColorPalette
    from pptx_generator.config.typography import Typography


class BaseTheme(ABC):
    """Abstract base class for presentation themes."""

    @property
    @abstractmethod
    def name(self) -> str:
        """Return the theme name."""
        pass

    @abstractmethod
    def apply_slide_background(self, slide: Slide, palette: ColorPalette) -> None:
        """Apply background styling to a slide.

        Args:
            slide: The slide to style
            palette: Color palette configuration
        """
        pass

    @abstractmethod
    def add_accent_bar(
        self, slide: Slide, palette: ColorPalette, position: str = "left"
    ) -> None:
        """Add a colored accent bar or stripe to the slide.

        Args:
            slide: The slide to modify
            palette: Color palette configuration
            position: Position of the accent bar ("left", "right", "top", "bottom")
        """
        pass

    @abstractmethod
    def add_footer(
        self,
        slide: Slide,
        palette: ColorPalette,
        typography: Typography,
        text: str,
        slide_number: int | None = None,
    ) -> None:
        """Add footer to the slide.

        Args:
            slide: The slide to modify
            palette: Color palette configuration
            typography: Typography configuration
            text: Footer text content
            slide_number: Optional slide number to display
        """
        pass

    @abstractmethod
    def style_shape(
        self, shape: BaseShape, palette: ColorPalette, style: str = "primary"
    ) -> None:
        """Apply styling to a shape.

        Args:
            shape: The shape to style
            palette: Color palette configuration
            style: Style variant ("primary", "secondary", "accent", "neutral")
        """
        pass

    @abstractmethod
    def add_decorative_element(
        self, slide: Slide, palette: ColorPalette, element_type: str = "corner"
    ) -> None:
        """Add decorative visual elements to the slide.

        Args:
            slide: The slide to modify
            palette: Color palette configuration
            element_type: Type of decoration ("corner", "geometric", "gradient", "none")
        """
        pass

    @staticmethod
    def _get_slide_dimensions(slide: Slide) -> tuple[int, int]:
        """Get slide width and height.

        Args:
            slide: The slide object

        Returns:
            Tuple of (width, height) in EMU units
        """
        # Access presentation dimensions through the slide's part
        presentation = slide.part.package.presentation_part.presentation
        return presentation.slide_width, presentation.slide_height

    @staticmethod
    def _hex_to_rgb(hex_color: str) -> RGBColor:
        """Convert hex color string to RGBColor.

        Args:
            hex_color: Hex color string (with or without #)

        Returns:
            RGBColor object
        """
        hex_color = hex_color.lstrip("#")
        return RGBColor(
            int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
        )

    @staticmethod
    def _set_shape_fill(shape: BaseShape, hex_color: str, transparency: float = 0.0) -> None:
        """Set solid fill color for a shape with optional transparency.

        Args:
            shape: The shape to fill
            hex_color: Hex color string
            transparency: Transparency level (0.0 = opaque, 1.0 = fully transparent)
        """
        shape.fill.solid()
        shape.fill.fore_color.rgb = BaseTheme._hex_to_rgb(hex_color)
        if transparency > 0.0:
            # Set transparency using the pptx-python API
            # Alpha is from 0 to 100000 (0% to 100%)
            alpha_value = int((1.0 - transparency) * 100000)
            try:
                from lxml import etree
                fill_elem = shape.fill._element
                srgb_clr = fill_elem.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr")
                if srgb_clr is not None:
                    alpha_elem = etree.SubElement(
                        srgb_clr,
                        "{http://schemas.openxmlformats.org/drawingml/2006/main}alpha"
                    )
                    alpha_elem.set("val", str(alpha_value))
            except (ImportError, AttributeError):
                # If lxml not available or element access fails, skip transparency
                pass

    @staticmethod
    def _set_shape_line(
        shape: BaseShape, hex_color: str, width_pt: float = 1.0
    ) -> None:
        """Set line/border styling for a shape.

        Args:
            shape: The shape to modify
            hex_color: Hex color string for the line
            width_pt: Line width in points
        """
        shape.line.color.rgb = BaseTheme._hex_to_rgb(hex_color)
        shape.line.width = Pt(width_pt)

    @staticmethod
    def _add_text_to_shape(
        shape: BaseShape,
        text: str,
        font_family: str,
        font_size_pt: int,
        font_color_hex: str,
        bold: bool = False,
        italic: bool = False,
    ) -> None:
        """Add styled text to a shape's text frame.

        Args:
            shape: The shape containing the text frame
            text: Text content
            font_family: Font family name
            font_size_pt: Font size in points
            font_color_hex: Font color as hex string
            bold: Whether to bold the text
            italic: Whether to italicize the text
        """
        text_frame = shape.text_frame
        text_frame.clear()
        paragraph = text_frame.paragraphs[0]
        paragraph.text = text

        for run in paragraph.runs:
            run.font.name = font_family
            run.font.size = Pt(font_size_pt)
            run.font.color.rgb = BaseTheme._hex_to_rgb(font_color_hex)
            run.font.bold = bold
            run.font.italic = italic
