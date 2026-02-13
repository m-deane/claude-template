"""Minimal theme - clean white background with subtle accents."""

from __future__ import annotations

from typing import TYPE_CHECKING

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from pptx_generator.themes.base import BaseTheme

if TYPE_CHECKING:
    from pptx.slide import Slide
    from pptx.shapes.base import BaseShape

    from pptx_generator.config.colors import ColorPalette
    from pptx_generator.config.typography import Typography


class MinimalTheme(BaseTheme):
    """Minimal theme with pure white background and subtle accents."""

    @property
    def name(self) -> str:
        """Return the theme name."""
        return "minimal"

    def apply_slide_background(self, slide: Slide, palette: ColorPalette) -> None:
        """Apply pure white background to the slide.

        Args:
            slide: The slide to style
            palette: Color palette configuration
        """
        background = slide.background
        fill = background.fill
        fill.solid()
        # Use white background
        fill.fore_color.rgb = self._hex_to_rgb("#FFFFFF")

    def add_accent_bar(
        self, slide: Slide, palette: ColorPalette, position: str = "top"
    ) -> None:
        """Add thin hairline accent on top edge only.

        Args:
            slide: The slide to modify
            palette: Color palette configuration
            position: Position of the accent bar (defaults to "top")
        """
        width, height = self._get_slide_dimensions(slide)

        if position == "top":
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0),
                Inches(0),
                width,
                Pt(2),
            )
        elif position == "bottom":
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0),
                height - Pt(2),
                width,
                Pt(2),
            )
        elif position == "left":
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0),
                Inches(0),
                Pt(2),
                height,
            )
        elif position == "right":
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                width - Pt(2),
                Inches(0),
                Pt(2),
                height,
            )
        else:
            # Default to top
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0),
                Inches(0),
                width,
                Pt(2),
            )

        self._set_shape_fill(bar, palette.accent)
        bar.line.fill.background()

    def add_footer(
        self,
        slide: Slide,
        palette: ColorPalette,
        typography: Typography,
        text: str,
        slide_number: int | None = None,
    ) -> None:
        """Add minimal footer with just page number, right-aligned.

        Args:
            slide: The slide to modify
            palette: Color palette configuration
            typography: Typography configuration
            text: Footer text content (optional, minimal theme prioritizes slide number)
            slide_number: Optional slide number to display
        """
        # Minimal theme shows only slide number, right-aligned
        if slide_number is not None:
            footer_text = str(slide_number)
        else:
            # If no slide number, show text but keep it minimal
            footer_text = text if text else ""

        if not footer_text:
            return

        width, height = self._get_slide_dimensions(slide)

        footer_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            width - Inches(1.0),
            height - Inches(0.4),
            Inches(0.8),
            Inches(0.3),
        )
        footer_box.fill.background()
        footer_box.line.fill.background()

        self._add_text_to_shape(
            footer_box,
            footer_text,
            typography.caption.family,
            typography.caption.size_pt - 2,  # Slightly smaller
            palette.neutral,
            bold=False,
            italic=False,
        )

        text_frame = footer_box.text_frame
        text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    def style_shape(
        self, shape: BaseShape, palette: ColorPalette, style: str = "primary"
    ) -> None:
        """Apply outline-only styling with no fills.

        Args:
            shape: The shape to style
            palette: Color palette configuration
            style: Style variant ("primary", "secondary", "accent", "neutral")
        """
        color_map = {
            "primary": palette.primary,
            "secondary": palette.secondary,
            "accent": palette.accent,
            "neutral": palette.neutral,
        }

        line_color = color_map.get(style, palette.primary)

        # No fill, outline only
        shape.fill.background()
        self._set_shape_line(shape, line_color, width_pt=1.5)

    def add_decorative_element(
        self, slide: Slide, palette: ColorPalette, element_type: str = "none"
    ) -> None:
        """Add no decorative elements (minimal theme principle).

        Args:
            slide: The slide to modify
            palette: Color palette configuration
            element_type: Type of decoration (always "none" for minimal theme)
        """
        # Minimal theme has no decorative elements
        pass
