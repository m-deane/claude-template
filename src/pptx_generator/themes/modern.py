"""Modern theme - bold, contemporary styling with gradients."""

from __future__ import annotations

from typing import TYPE_CHECKING

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from pptx_generator.themes.base import BaseTheme

if TYPE_CHECKING:
    from pptx.slide import Slide
    from pptx.shapes.base import BaseShape

    from pptx_generator.config.colors import ColorPalette
    from pptx_generator.config.typography import Typography


class ModernTheme(BaseTheme):
    """Modern theme with bold, contemporary styling."""

    @property
    def name(self) -> str:
        """Return the theme name."""
        return "modern"

    def apply_slide_background(self, slide: Slide, palette: ColorPalette) -> None:
        """Apply background with full-width colored header band.

        Args:
            slide: The slide to style
            palette: Color palette configuration
        """
        # Set base background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self._hex_to_rgb(palette.background)

        # Add colored header band (top 15% of slide)
        width, height = self._get_slide_dimensions(slide)
        header_band = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            width,
            int(height * 0.15),
        )
        self._set_shape_fill(header_band, palette.primary)
        header_band.line.fill.background()

    def add_accent_bar(
        self, slide: Slide, palette: ColorPalette, position: str = "left"
    ) -> None:
        """Add wide accent bar on left edge.

        Args:
            slide: The slide to modify
            palette: Color palette configuration
            position: Position of the accent bar (defaults to "left")
        """
        width, height = self._get_slide_dimensions(slide)

        if position == "left":
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0),
                Inches(0),
                Inches(0.35),
                height,
            )
            fill_color = palette.accent
        elif position == "right":
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                width - Inches(0.35),
                Inches(0),
                Inches(0.35),
                height,
            )
            fill_color = palette.accent
        elif position == "top":
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0),
                Inches(0),
                width,
                Inches(0.35),
            )
            fill_color = palette.secondary
        elif position == "bottom":
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0),
                height - Inches(0.35),
                width,
                Inches(0.35),
            )
            fill_color = palette.secondary
        else:
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0),
                Inches(0),
                Inches(0.35),
                height,
            )
            fill_color = palette.accent

        self._set_shape_fill(bar, fill_color)
        bar.line.fill.background()

    def add_footer(
        self,
        slide: Slide,
        palette: ColorPalette,
        typography: Typography,
        text: str,
        slide_number: int | None = None,
    ) -> None:
        """Add modern footer with bold styling.

        Args:
            slide: The slide to modify
            palette: Color palette configuration
            typography: Typography configuration
            text: Footer text content
            slide_number: Optional slide number to display
        """
        width, height = self._get_slide_dimensions(slide)

        footer_text = text
        if slide_number is not None:
            footer_text = f"{text}  •  {slide_number}"

        footer_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5),
            height - Inches(0.45),
            width - Inches(1.0),
            Inches(0.35),
        )
        footer_box.fill.background()
        footer_box.line.fill.background()

        self._add_text_to_shape(
            footer_box,
            footer_text,
            typography.caption.family,
            typography.caption.size_pt,
            palette.accent,
            bold=True,
            italic=typography.caption.italic,
        )

        text_frame = footer_box.text_frame
        text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    def style_shape(
        self, shape: BaseShape, palette: ColorPalette, style: str = "primary"
    ) -> None:
        """Apply gradient fill styling to shapes.

        Args:
            shape: The shape to style
            palette: Color palette configuration
            style: Style variant ("primary", "secondary", "accent", "neutral")
        """
        color_map = {
            "primary": (palette.gradient_start, palette.gradient_end),
            "secondary": (palette.secondary, palette.accent),
            "accent": (palette.accent, palette.primary),
            "neutral": (palette.neutral, palette.background),
        }

        start_color, end_color = color_map.get(style, (palette.gradient_start, palette.gradient_end))

        # Apply gradient fill
        fill = shape.fill
        fill.gradient()
        fill.gradient_angle = 45.0
        fill.gradient_stops[0].color.rgb = self._hex_to_rgb(start_color)
        fill.gradient_stops[1].color.rgb = self._hex_to_rgb(end_color)

        # Add subtle border
        self._set_shape_line(shape, start_color, width_pt=2.0)

    def add_decorative_element(
        self, slide: Slide, palette: ColorPalette, element_type: str = "geometric"
    ) -> None:
        """Add geometric decorative elements (overlapping circles/rectangles).

        Args:
            slide: The slide to modify
            palette: Color palette configuration
            element_type: Type of decoration (defaults to "geometric")
        """
        if element_type == "geometric":
            width, height = self._get_slide_dimensions(slide)

            # Add overlapping circles in bottom-right
            circle1 = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                width - Inches(2.0),
                height - Inches(1.8),
                Inches(1.2),
                Inches(1.2),
            )
            self._set_shape_fill(circle1, palette.accent, transparency=0.6)
            circle1.line.fill.background()

            circle2 = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                width - Inches(1.5),
                height - Inches(1.5),
                Inches(1.0),
                Inches(1.0),
            )
            self._set_shape_fill(circle2, palette.secondary, transparency=0.5)
            circle2.line.fill.background()

            # Add rectangle overlay
            rect = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                width - Inches(1.8),
                height - Inches(1.2),
                Inches(0.8),
                Inches(0.8),
            )
            self._set_shape_fill(rect, palette.primary, transparency=0.7)
            rect.line.fill.background()
            rect.rotation = 15
        elif element_type == "none":
            pass
        else:
            # Default to geometric decoration
            self.add_decorative_element(slide, palette, "geometric")
