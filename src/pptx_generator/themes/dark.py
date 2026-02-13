"""Dark theme - charcoal background with glowing accents."""

from __future__ import annotations

from typing import TYPE_CHECKING

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from pptx_generator.themes.base import BaseTheme

if TYPE_CHECKING:
    from pptx.slide import Slide
    from pptx.shapes.base import BaseShape

    from pptx_generator.config.colors import ColorPalette
    from pptx_generator.config.typography import Typography


class DarkTheme(BaseTheme):
    """Dark theme with charcoal background and glowing accents."""

    @property
    def name(self) -> str:
        """Return the theme name."""
        return "dark"

    def apply_slide_background(self, slide: Slide, palette: ColorPalette) -> None:
        """Apply dark charcoal background to the slide.

        Args:
            slide: The slide to style
            palette: Color palette configuration
        """
        background = slide.background
        fill = background.fill
        fill.solid()
        # Use dark background color, or fallback to near-black
        bg_color = palette.background if palette.background.startswith("#") else "#1a1a1a"
        fill.fore_color.rgb = self._hex_to_rgb(bg_color)

    def add_accent_bar(
        self, slide: Slide, palette: ColorPalette, position: str = "left"
    ) -> None:
        """Add glowing accent lines to the slide.

        Args:
            slide: The slide to modify
            palette: Color palette configuration
            position: Position of the accent bar
        """
        width, height = self._get_slide_dimensions(slide)

        if position == "left":
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0),
                Inches(0),
                Inches(0.08),
                height,
            )
        elif position == "right":
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                width - Inches(0.08),
                Inches(0),
                Inches(0.08),
                height,
            )
        elif position == "top":
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0),
                Inches(0),
                width,
                Inches(0.08),
            )
        elif position == "bottom":
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0),
                height - Inches(0.08),
                width,
                Inches(0.08),
            )
        else:
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0),
                Inches(0),
                Inches(0.08),
                height,
            )

        # Bright glowing accent color
        self._set_shape_fill(bar, palette.accent)
        bar.line.fill.background()

    def add_footer(
        self,
        slide: Slide,
        palette: ColorPalette,
        typography: Typography,
        text: str,
        slide_number: int | None = None,
    ) -> None:
        """Add subtle footer with gradient bar at bottom.

        Args:
            slide: The slide to modify
            palette: Color palette configuration
            typography: Typography configuration
            text: Footer text content
            slide_number: Optional slide number to display
        """
        # Add subtle bottom gradient bar instead of corner decoration
        width, height = self._get_slide_dimensions(slide)
        gradient_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            height - Inches(0.25),
            width,
            Inches(0.25),
        )

        fill = gradient_bar.fill
        fill.gradient()
        fill.gradient_angle = 0.0
        fill.gradient_stops[0].color.rgb = self._hex_to_rgb(palette.accent)
        fill.gradient_stops[1].color.rgb = self._hex_to_rgb(palette.background)
        gradient_bar.line.fill.background()

        # Add footer text
        footer_text = text
        if slide_number is not None:
            footer_text = f"{text}  |  {slide_number}"

        footer_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5),
            height - Inches(0.5),
            width - Inches(1.0),
            Inches(0.35),
        )
        footer_box.fill.background()
        footer_box.line.fill.background()

        self._add_text_to_shape(
            footer_box,
            footer_text,
            typography.caption.family,
            typography.caption.size_pt,
            typography.caption.color_hex or palette.text_light,
            bold=typography.caption.bold,
            italic=typography.caption.italic,
        )

        text_frame = footer_box.text_frame
        text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    def style_shape(
        self, shape: BaseShape, palette: ColorPalette, style: str = "primary"
    ) -> None:
        """Apply semi-transparent fills with bright borders.

        Args:
            shape: The shape to style
            palette: Color palette configuration
            style: Style variant ("primary", "secondary", "accent", "neutral")
        """
        color_map = {
            "primary": palette.primary,
            "secondary": palette.secondary,
            "accent": palette.accent,
            "neutral": palette.neutral,
        }

        fill_color = color_map.get(style, palette.primary)

        # Semi-transparent fill
        self._set_shape_fill(shape, fill_color, transparency=0.4)

        # Bright glowing border
        self._set_shape_line(shape, palette.accent, width_pt=3.0)

    def add_decorative_element(
        self, slide: Slide, palette: ColorPalette, element_type: str = "gradient"
    ) -> None:
        """Add subtle gradient decoration (handled in footer for dark theme).

        Args:
            slide: The slide to modify
            palette: Color palette configuration
            element_type: Type of decoration (defaults to "gradient")
        """
        # Dark theme uses gradient bar in footer instead of separate decoration
        # This method is intentionally minimal to avoid visual clutter
        if element_type == "gradient":
            # Gradient bar is added as part of the footer
            pass
        elif element_type == "none":
            pass
        else:
            # Default: no additional decoration
            pass
